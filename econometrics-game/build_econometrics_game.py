"""Build the CRUK Econometrics Game workbook from Analytic Partners' Wave 1 deck.

Usage: python build_game.py <wave1_results.pptx> <output.xlsx>

Every weekly number in the workbook is read from the charts in the deck. The workbook
has no macros. All game logic is Excel formulas.
"""
import sys, datetime, math
import numpy as np
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter, column_index_from_string
from openpyxl.chart import LineChart, BarChart, Reference, Series
from openpyxl.chart.series import SeriesLabel
from openpyxl.chart.shapes import GraphicalProperties
from openpyxl.chart.layout import Layout, ManualLayout
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.formatting.rule import FormulaRule, CellIsRule
from openpyxl.comments import Comment

PPTX = sys.argv[1]
OUT = sys.argv[2]

# ----------------------------------------------------------------------------------
# 1. Read AP's weekly series from the deck charts
# ----------------------------------------------------------------------------------
prs = Presentation(PPTX)
series = {}
for idx in [16, 21, 23, 35, 46]:
    s = prs.slides[idx - 1]
    for sh in s.shapes:
        if not sh.has_chart:
            continue
        for pi, plot in enumerate(sh.chart.plots):
            cats = list(plot.categories)
            for ser in plot.series:
                vals = np.array([0.0 if v is None else float(v) for v in ser.values])
                series[(idx, ser.name.strip())] = vals
            if idx == 23:
                dates = [datetime.date(1899, 12, 30) + datetime.timedelta(days=int(float(c))) for c in cats]

N = len(dates)
assert N == 148, N
yrs = np.array([d.year for d in dates])
mons = np.array([d.month for d in dates])

def S(idx, name):
    return series[(idx, name)]

actual = S(16, 'RFL Sign Ups')
cog_signups = S(16, 'COG Sign Ups')
consideration = S(46, 'Brand Consideration')
total_media = S(46, 'Total Media Spend')

dec = {name: v for (i, name), v in series.items() if i == 23}
sp = {name: v for (i, name), v in series.items() if i == 21}

# AP blocks (real weekly contributions from the Wave 1 decomposition)
ap_base = dec['Brand (Consideration)'] + dec['Emails']
ap_brake = dec['Price & Distribution'] + dec['Seasonality']           # negative
ap_sale = dec['Discount']
rfl_channel_contrib = ['BRTV', 'DRTV', 'VOD', 'OOH', 'Radio', 'Regional', 'Paid Social', 'Paid Search',
                       'Digital Audio', 'Display', 'YouTube', 'Press Partnership', 'Other Media']
ap_rfl = sum(dec[k] for k in rfl_channel_contrib)
ap_cog = dec['Halo (COG)']
ap_leg = dec['Halo (Legacy)']
ap_brand = dec['Brand (TV, VOD, OOH)']
channel_groups = [
    ('TV (BRTV and DRTV)', ['BRTV RFL', 'DRTV RFL'], ['BRTV', 'DRTV']),
    ('Online video (VOD)', ['VOD'], ['VOD']),
    ('Radio', ['Radio'], ['Radio']),
    ('Posters (OOH)', ['OOH'], ['OOH']),
    ('Digital, social and the rest', ['Paid Social EMC', 'Paid Social Oliver', 'Google Ads Generic',
                                       'Microsoft Ads Generic', 'Display', 'Digital Audio', 'YouTube', 'Regional',
                                       'Press Partnership', 'Telemarketing Ethicall', 'Telemarketing SS',
                                       'Direct Mail', 'Door Drops'],
     ['Paid Social', 'Paid Search', 'Display', 'Digital Audio', 'YouTube', 'Regional', 'Press Partnership',
      'Other Media']),
]
ap_channels = [sum(dec[c] for c in ck) for _, _, ck in channel_groups]
ch_spend = [sum(sp[s] for s in sk) / 1000 for _, sk, _ in channel_groups]   # £k
rfl_spend = sum(sp[k] for k in sp if k not in ('BRTV Brand', 'VOD Brand')) / 1000
cog_spend = S(16, 'COG') / 1000
leg_spend = S(16, 'Legacy') / 1000
brand_spend = S(16, 'Brand') / 1000
ap_total = sum(dec.values())

# shapes for the non-media blocks (dial x shape = block; the answer dial reproduces AP exactly)
ans_base = ap_base.mean()
ans_brake = -ap_brake.mean()
ans_sale = ap_sale[ap_sale > 0].mean()
shape_base = ap_base / ans_base
shape_brake = -ap_brake / ans_brake
shape_sale = ap_sale / ans_sale


def adstock(x, m):
    out = np.zeros_like(x)
    a = 0.0
    for i, v in enumerate(x):
        a = v + m * a
        out[i] = a
    return out


def fit_block(contrib, spend_k):
    """dial = sign-ups per £1,000 counting every week the effect lasts; memory in tenths."""
    best = None
    for m10 in range(0, 10):
        m = m10 / 10
        A = (1 - m) * adstock(spend_k, m)
        k = (A * contrib).sum() / (A * A).sum()
        sse = ((contrib - k * A) ** 2).sum()
        if best is None or sse < best[2]:
            best = (m10 * 10, k, sse)
    return best[0], float(contrib.sum() / spend_k.sum())


mem_rfl, ans_rfl = fit_block(ap_rfl, rfl_spend)
mem_cog, ans_cog = fit_block(ap_cog, cog_spend)
mem_leg, ans_leg = fit_block(ap_leg, leg_spend)
mem_brand, ans_brand = fit_block(ap_brand, brand_spend)
ch_fits = [fit_block(c, s) for c, s in zip(ap_channels, ch_spend)]
print('answers: base %.0f brake %.0f sale %.0f | rfl %.1f mem %d | cog %.2f mem %d | leg %.2f mem %d | brand %.2f mem %d'
      % (ans_base, ans_brake, ans_sale, ans_rfl, mem_rfl, ans_cog, mem_cog, ans_leg, mem_leg, ans_brand, mem_brand))
for (g, _, _), (m, k) in zip(channel_groups, ch_fits):
    print('   %-32s %.2f per £k, memory %d' % (g, k, m))


def r2(model):
    return 1 - ((actual - model) ** 2).sum() / ((actual - actual.mean()) ** 2).sum()


# answer model, cumulative by level, to set targets
b1 = ans_base * shape_base
b2 = -ans_brake * shape_brake
b3 = ans_sale * shape_sale
b4 = ans_rfl * (1 - mem_rfl / 100) * adstock(rfl_spend, mem_rfl / 100)
b5 = (ans_cog * (1 - mem_cog / 100) * adstock(cog_spend, mem_cog / 100)
      + ans_leg * (1 - mem_leg / 100) * adstock(leg_spend, mem_leg / 100)
      + ans_brand * (1 - mem_brand / 100) * adstock(brand_spend, mem_brand / 100))
b6 = sum(k * (1 - m / 100) * adstock(s, m / 100) for (m, k), s in zip(ch_fits, ch_spend))
m3 = b1 + b2 + b3
m4 = m3 + b4
m5 = m4 + b5
m6 = m3 + b6 + b5
targets = {3: r2(m3), 4: r2(m4), 5: r2(m5), 6: r2(m6)}
print('answer R2 by level', {k: round(v, 3) for k, v in targets.items()}, 'AP decomposition R2', round(r2(ap_total), 3))
ans_r2 = r2(m6)
resid = actual - m6
ans_dw = (np.diff(resid) ** 2).sum() / (resid ** 2).sum()
mask = actual > 2000
ans_mape = (np.abs(actual - m6)[mask] / actual[mask]).mean()
print('answer DW %.2f MAPE %.1f%%' % (ans_dw, ans_mape * 100))

# ----------------------------------------------------------------------------------
# 2. Styles
# ----------------------------------------------------------------------------------
NAVY = '1B2A4A'
STEEL = '4A6FA5'
MAGENTA = 'E60078'
CYAN = '009CEE'
GREY = '878787'
LIGHTBLUE = 'DDEBF7'
GREEN = '2E7D32'
GREENFILL = 'E6F4EA'
AMBERFILL = 'FFF4E5'
AMBER = 'C75000'
PALE = 'F4F6FA'

F_TITLE = Font(name='Arial', size=16, bold=True, color=NAVY)
F_SUB = Font(name='Arial', size=11, italic=True, color='444444')
F_H = Font(name='Arial', size=12, bold=True, color=STEEL)
F_BODY = Font(name='Arial', size=11, color='1A1A1A')
F_BOLD = Font(name='Arial', size=11, bold=True, color='1A1A1A')
F_SMALL = Font(name='Arial', size=9, italic=True, color='444444')
F_INPUT = Font(name='Arial', size=12, bold=True, color='0000FF')
F_BIG = Font(name='Arial', size=14, bold=True, color=NAVY)
F_WHITE = Font(name='Arial', size=11, bold=True, color='FFFFFF')
FILL_INPUT = PatternFill('solid', fgColor=LIGHTBLUE)
FILL_HEAD = PatternFill('solid', fgColor=NAVY)
FILL_PALE = PatternFill('solid', fgColor=PALE)
FILL_ACCENT = PatternFill('solid', fgColor='E8F0F8')
WRAP = Alignment(wrap_text=True, vertical='top')
WRAP_C = Alignment(wrap_text=True, vertical='center')
CENTER = Alignment(horizontal='center', vertical='center', wrap_text=True)
thin = Side(style='thin', color='BBC7D9')
BOX = Border(left=thin, right=thin, top=thin, bottom=thin)
ACC_SIDE = Side(style='medium', color=STEEL)

wb = Workbook()
wb.remove(wb.active)

ENG = 'Under the bonnet'
E = "'" + ENG + "'!"
ANS = 'Answers'
A_ = ANS + '!'
R0 = 6              # first data row in engine
R1 = R0 + N - 1     # last data row


def rng(col):
    return f"{E}${col}${R0}:${col}${R1}"


def set_widths(ws, widths):
    for col, w in widths.items():
        ws.column_dimensions[col].width = w


def text_row(ws, row, text, font=F_BODY, merge_to='D', height=None, fill=None, indent=0):
    c = ws.cell(row=row, column=2, value=text)
    c.font = font
    c.alignment = Alignment(wrap_text=True, vertical='top', indent=indent)
    if merge_to:
        ws.merge_cells(f'B{row}:{merge_to}{row}')
    if fill:
        for col in 'BCD':
            ws[f'{col}{row}'].fill = fill
    if height is None:
        if merge_to:
            last = column_index_from_string(merge_to)
            width_units = sum((ws.column_dimensions[get_column_letter(k)].width or 8.43) for k in range(2, last + 1))
        else:
            width_units = ws.column_dimensions['B'].width or 8.43
        width_chars = max(20, int(width_units * 1.05))
        lines = max(1, math.ceil(len(text) / width_chars) + text.count('\n'))
        height = 15 * lines + 4
    ws.row_dimensions[row].height = height
    return c


def header(ws, row, text):
    c = ws.cell(row=row, column=2, value=text)
    c.font = F_H
    ws.row_dimensions[row].height = 20


def title(ws, text, sub):
    ws['B1'] = text
    ws['B1'].font = F_TITLE
    ws.row_dimensions[1].height = 30
    ws['B2'] = sub
    ws['B2'].font = F_SUB
    ws.merge_cells('B2:N2')
    ws.row_dimensions[2].height = 18
    ws.sheet_view.showGridLines = False


def input_cell(ws, ref, value=None, fmt='#,##0'):
    c = ws[ref]
    if value is not None:
        c.value = value
    c.font = F_INPUT
    c.fill = FILL_INPUT
    c.border = BOX
    c.alignment = Alignment(horizontal='center', vertical='center')
    c.number_format = fmt
    return c


def label(ws, ref, text, bold=False, wrap=True):
    c = ws[ref]
    c.value = text
    c.font = F_BOLD if bold else F_BODY
    c.alignment = Alignment(wrap_text=wrap, vertical='center')
    return c


def hint(ws, ref, text):
    c = ws[ref]
    c.value = text
    c.font = F_SMALL
    c.alignment = Alignment(wrap_text=True, vertical='center')
    return c


def verdict_formats(ws, cell_range):
    ws.conditional_formatting.add(cell_range, FormulaRule(
        formula=[f'ISNUMBER(SEARCH("About right",{cell_range.split(":")[0]}))'],
        fill=PatternFill('solid', fgColor=GREENFILL), font=Font(name='Arial', bold=True, color=GREEN)))
    ws.conditional_formatting.add(cell_range, FormulaRule(
        formula=[f'ISNUMBER(SEARCH("Too",{cell_range.split(":")[0]}))'],
        fill=PatternFill('solid', fgColor=AMBERFILL), font=Font(name='Arial', bold=True, color=AMBER)))
    ws.conditional_formatting.add(cell_range, FormulaRule(
        formula=[f'ISNUMBER(SEARCH("done",{cell_range.split(":")[0]}))'],
        fill=PatternFill('solid', fgColor=GREENFILL), font=Font(name='Arial', bold=True, color=GREEN)))
    ws.conditional_formatting.add(cell_range, FormulaRule(
        formula=[f'ISNUMBER(SEARCH("Keep going",{cell_range.split(":")[0]}))'],
        fill=PatternFill('solid', fgColor=AMBERFILL), font=Font(name='Arial', bold=True, color=AMBER)))


def yes_no(ws, ref):
    dv = DataValidation(type='list', formula1='"No,Yes"', allow_blank=False)
    ws.add_data_validation(dv)
    dv.add(ref)
    c = ws[ref]
    c.value = 'No'
    c.font = F_INPUT
    c.fill = FILL_INPUT
    c.border = BOX
    c.alignment = Alignment(horizontal='center')


def memory_dropdown(ws, ref):
    dv = DataValidation(type='list', formula1='"0,10,20,30,40,50,60,70,80,90"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add(ref)
    input_cell(ws, ref, None, fmt='0"%"')


def name_series(s, text):
    s.tx = SeriesLabel(v=text)


def line_style(s, color, width_pt=1.75, dash=None, smooth=False):
    s.graphicalProperties.line.solidFill = color
    s.graphicalProperties.line.width = int(width_pt * 12700)
    if dash:
        s.graphicalProperties.line.dashStyle = dash
    s.smooth = smooth
    s.marker.symbol = 'none'


def base_chart(ch, title_text, y_title='Sign-ups a week', width=24, height=11):
    ch.title = title_text
    ch.style = 2
    ch.width = width
    ch.height = height
    ch.legend.position = 'b'
    ch.y_axis.title = y_title
    ch.y_axis.number_format = '#,##0'
    ch.y_axis.majorGridlines = None
    ch.x_axis.number_format = 'mmm yy'
    ch.x_axis.delete = False
    ch.y_axis.delete = False
    try:
        ch.x_axis.tickLblSkip = 13
        ch.x_axis.tickMarkSkip = 13
    except Exception:
        pass


def model_chart(ws, anchor, level, cols, title_text, width=24, height=11):
    """Actual (magenta), answer line (grey dashed, hidden under yours until revealed), your model (cyan)."""
    ch = LineChart()
    base_chart(ch, title_text, width=width, height=height)
    cats = Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1)
    for col, name, color, dash in cols:
        ref = Reference(wb[ENG], min_col=col, min_row=R0 - 1, max_row=R1)
        ch.add_data(ref, titles_from_data=True)
        s = ch.series[-1]
        name_series(s, name)
        line_style(s, color, dash=dash)
    ch.set_categories(cats)
    ws.add_chart(ch, anchor)
    return ch


# ----------------------------------------------------------------------------------
# 3. Engine sheet ("Under the bonnet")
# ----------------------------------------------------------------------------------
eng = wb.create_sheet(ENG)
eng.sheet_properties.tabColor = GREY
eng['A1'] = 'Under the bonnet: every weekly number and every formula the game uses. Nothing is hidden.'
eng['A1'].font = F_TITLE
eng['A2'] = ("Columns A to D and the AP columns are Analytic Partners' Wave 1 weekly figures, read from the charts in "
             "their July 2026 results deck. Everything else is a formula that points at the blue cells on the level tabs. "
             "Click any cell to see how it is built. Spend is in thousands of pounds.")
eng['A2'].font = F_BODY
eng.merge_cells('A2:T2')
eng.row_dimensions[2].height = 34
eng['A3'] = ("Shape columns: AP's weekly line for that block divided by its typical week, so that your dial times the "
             "shape gives your block. Adstock: this week's spend plus memory times last week's total.")
eng['A3'].font = F_SMALL
eng.merge_cells('A3:T3')
eng.freeze_panes = 'E6'

# level sheet cell map (fixed positions used by engine formulas)
L = {i: f"'Level {i}'!" for i in range(1, 10)}
D = {  # dial cells
    'base': L[1] + '$C$11', 'brake': L[2] + '$C$11', 'sale': L[3] + '$C$11',
    'rfl': L[4] + '$C$11', 'rfl_mem': L[4] + '$C$12',
    'cog': L[5] + '$C$11', 'leg': L[5] + '$C$12', 'brand': L[5] + '$C$13',
    'ch': [L[6] + f'$C${11 + i}' for i in range(5)],
}
REVEAL = {i: L[i] + '$C$31' for i in range(1, 8)}
POUND = "'Start here'!$C$38"

cols = {}
hdrs = [
    ('A', 'Week starting'), ('B', 'Year'), ('C', 'Month'), ('D', 'Real sign-ups'),
    ('E', 'Shape: regulars'), ('F', 'Shape: brakes'), ('G', 'Shape: sale'),
    ('H', 'RFL adverts £k'), ('I', 'RFL adstock'),
    ('J', 'Committed Giving adverts £k'), ('K', 'Legacy adverts £k'), ('L', 'Brand adverts £k'), ('M', 'Brand adstock'),
    ('N', 'TV £k'), ('O', 'Online video £k'), ('P', 'Radio £k'), ('Q', 'Posters £k'), ('R', 'Digital and rest £k'),
    ('S', 'TV adstock'), ('T', 'Online video adstock'), ('U', 'Radio adstock'), ('V', 'Posters adstock'), ('W', 'Digital adstock'),
    ('X', 'Your block 1: regulars'), ('Y', 'Your block 2: brakes'), ('Z', 'Your block 3: sale'),
    ('AA', 'Your block 4: RFL adverts'), ('AB', 'Your block 5: CG halo'), ('AC', 'Your block 5: Legacy halo'),
    ('AD', 'Your block 5: Brand adverts'),
    ('AE', 'Your block 6: TV'), ('AF', 'Your block 6: Online video'), ('AG', 'Your block 6: Radio'),
    ('AH', 'Your block 6: Posters'), ('AI', 'Your block 6: Digital and rest'),
    ('AJ', 'Your model after level 1'), ('AK', 'Your model after level 2'), ('AL', 'Your model after level 3'),
    ('AM', 'Your model after level 4'), ('AN', 'Your model after level 5'), ('AO', 'Your model after level 6'),
    ('AP', 'Your final model'),
    ('AQ', 'AP block: regulars'), ('AR', 'AP block: brakes'), ('AS', 'AP block: sale'), ('AT', 'AP block: RFL adverts'),
    ('AU', 'AP block: CG halo'), ('AV', 'AP block: Legacy halo'), ('AW', 'AP block: Brand adverts'),
    ('AX', 'AP block: TV'), ('AY', 'AP block: Online video'), ('AZ', 'AP block: Radio'), ('BA', 'AP block: Posters'),
    ('BB', 'AP block: Digital and rest'), ('BC', 'Wave 1 model'),
    ('BD', 'Wave 1 after level 1'), ('BE', 'Wave 1 after level 2'), ('BF', 'Wave 1 after level 3'),
    ('BG', 'Wave 1 after level 4'), ('BH', 'Wave 1 after level 5'), ('BI', 'Wave 1 after level 6'),
    ('BJ', 'Answer line L1'), ('BK', 'Answer line L2'), ('BL', 'Answer line L3'), ('BM', 'Answer line L4'),
    ('BN', 'Answer line L5'), ('BO', 'Answer line L6'), ('BP', 'Answer line L7'),
    ('BQ', 'Miss (real minus yours)'), ('BR', 'Miss squared'), ('BS', 'Miss change squared'), ('BT', 'Weekly % miss (weeks over 2,000)'),
    ('BU', 'Real new Committed Givers'), ('BV', 'Brand consideration'), ('BW', 'All CRUK media £k'),
    ('BX', 'Your block 4 on this level 4 chart'), ('BY', 'RFL adverts £k (for chart)'),
    ('BZ', 'Sale weeks (for chart)'), ('CA', 'Your halo blocks'), ('CB', 'Your split blocks'),
]
for col, h in hdrs:
    c = eng[f'{col}5']
    c.value = h
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
eng.row_dimensions[5].height = 44
for col, _ in hdrs:
    eng.column_dimensions[col].width = 13
eng.column_dimensions['A'].width = 12

for i in range(N):
    r = R0 + i
    p = r - 1
    eng[f'A{r}'] = dates[i]
    eng[f'A{r}'].number_format = 'dd mmm yy'
    eng[f'B{r}'] = f'=YEAR(A{r})'
    eng[f'C{r}'] = f'=MONTH(A{r})'
    eng[f'D{r}'] = round(float(actual[i]), 1)
    eng[f'E{r}'] = round(float(shape_base[i]), 4)
    eng[f'F{r}'] = round(float(shape_brake[i]), 4)
    eng[f'G{r}'] = round(float(shape_sale[i]), 4)
    eng[f'H{r}'] = round(float(rfl_spend[i]), 2)
    eng[f'I{r}'] = f'=H{r}' if i == 0 else f'=H{r}+{D["rfl_mem"]}/100*I{p}'
    eng[f'J{r}'] = round(float(cog_spend[i]), 2)
    eng[f'K{r}'] = round(float(leg_spend[i]), 2)
    eng[f'L{r}'] = round(float(brand_spend[i]), 2)
    eng[f'M{r}'] = f'=L{r}' if i == 0 else f'=L{r}+{A_}$C$14/100*M{p}'
    for j, colx in enumerate('NOPQR'):
        eng[f'{colx}{r}'] = round(float(ch_spend[j][i]), 2)
    for j, (colx, spcol) in enumerate(zip('STUVW', 'NOPQR')):
        eng[f'{colx}{r}'] = f'={spcol}{r}' if i == 0 else f'={spcol}{r}+{A_}$C${16 + j}/100*{colx}{p}'
    eng[f'X{r}'] = f'=IF({D["base"]}="",0,{D["base"]})*E{r}'
    eng[f'Y{r}'] = f'=-IF({D["brake"]}="",0,{D["brake"]})*F{r}'
    eng[f'Z{r}'] = f'=IF({D["sale"]}="",0,{D["sale"]})*G{r}'
    eng[f'AA{r}'] = f'=IF({D["rfl"]}="",0,{D["rfl"]})*(1-{D["rfl_mem"]}/100)*I{r}'
    eng[f'AB{r}'] = f'=IF({D["cog"]}="",0,{D["cog"]})*J{r}'
    eng[f'AC{r}'] = f'=IF({D["leg"]}="",0,{D["leg"]})*K{r}'
    eng[f'AD{r}'] = f'=IF({D["brand"]}="",0,{D["brand"]})*(1-{A_}$C$14/100)*M{r}'
    for j, (colx, adcol) in enumerate(zip(['AE', 'AF', 'AG', 'AH', 'AI'], 'STUVW')):
        eng[f'{colx}{r}'] = f'=IF({D["ch"][j]}="",0,{D["ch"][j]})*(1-{A_}$C${16 + j}/100)*{adcol}{r}'
    eng[f'AJ{r}'] = f'=X{r}'
    eng[f'AK{r}'] = f'=X{r}+Y{r}'
    eng[f'AL{r}'] = f'=AK{r}+Z{r}'
    eng[f'AM{r}'] = f'=AL{r}+AA{r}'
    eng[f'AN{r}'] = f'=AM{r}+AB{r}+AC{r}+AD{r}'
    eng[f'AO{r}'] = f'=AL{r}+AE{r}+AF{r}+AG{r}+AH{r}+AI{r}+AB{r}+AC{r}+AD{r}'
    eng[f'AP{r}'] = f'=IF({A_}$C$22="Yes",AO{r},AN{r})'
    eng[f'AQ{r}'] = round(float(ap_base[i]), 1)
    eng[f'AR{r}'] = round(float(ap_brake[i]), 1)
    eng[f'AS{r}'] = round(float(ap_sale[i]), 1)
    eng[f'AT{r}'] = round(float(ap_rfl[i]), 1)
    eng[f'AU{r}'] = round(float(ap_cog[i]), 1)
    eng[f'AV{r}'] = round(float(ap_leg[i]), 1)
    eng[f'AW{r}'] = round(float(ap_brand[i]), 1)
    for j, colx in enumerate(['AX', 'AY', 'AZ', 'BA', 'BB']):
        eng[f'{colx}{r}'] = round(float(ap_channels[j][i]), 1)
    eng[f'BC{r}'] = f'=AQ{r}+AR{r}+AS{r}+AT{r}+AU{r}+AV{r}+AW{r}'
    eng[f'BD{r}'] = f'=AQ{r}'
    eng[f'BE{r}'] = f'=AQ{r}+AR{r}'
    eng[f'BF{r}'] = f'=BE{r}+AS{r}'
    eng[f'BG{r}'] = f'=BF{r}+AT{r}'
    eng[f'BH{r}'] = f'=BG{r}+AU{r}+AV{r}+AW{r}'
    eng[f'BI{r}'] = f'=BC{r}'
    for j, (colx, apc, mine) in enumerate(zip(['BJ', 'BK', 'BL', 'BM', 'BN', 'BO', 'BP'],
                                              ['BD', 'BE', 'BF', 'BG', 'BH', 'BI', 'BC'],
                                              ['AJ', 'AK', 'AL', 'AM', 'AN', 'AO', 'AP'])):
        eng[f'{colx}{r}'] = f'=IF({REVEAL[j + 1]}="Yes",{apc}{r},{mine}{r})'
    eng[f'BQ{r}'] = f'=D{r}-AP{r}'
    eng[f'BR{r}'] = f'=BQ{r}^2'
    eng[f'BS{r}'] = 0 if i == 0 else f'=(BQ{r}-BQ{p})^2'
    eng[f'BT{r}'] = f'=IF(D{r}>2000,ABS(BQ{r})/D{r},"")'
    eng[f'BU{r}'] = round(float(cog_signups[i]), 0)
    eng[f'BV{r}'] = round(float(consideration[i]), 4)
    eng[f'BV{r}'].number_format = '0.0%'
    eng[f'BW{r}'] = round(float(total_media[i]) / 1000, 1)
    eng[f'BX{r}'] = f'=AA{r}'
    eng[f'BY{r}'] = f'=H{r}'
    eng[f'BZ{r}'] = f'=Z{r}'
    eng[f'CA{r}'] = f'=AB{r}+AC{r}+AD{r}'
    eng[f'CB{r}'] = f'=AE{r}+AF{r}+AG{r}+AH{r}+AI{r}'
    for col, _ in hdrs:
        cc = eng[f'{col}{r}']
        cc.font = Font(name='Arial', size=9)
        if col not in ('A', 'B', 'C', 'BV'):
            cc.number_format = '#,##0.0' if col in 'EFG' else '#,##0'

# named engine columns for formulas
CL = {h: col for col, h in hdrs}
C = {h: column_index_from_string(col) for col, h in hdrs}
ACT = rng('D')
FINAL = rng('AP')
YEAR = rng('B')
MONTH = rng('C')

# ----------------------------------------------------------------------------------
# 4. Answers sheet (very hidden)
# ----------------------------------------------------------------------------------
ans = wb.create_sheet(ANS)
ans['A1'] = 'Answer key. Hidden so that the game stays a game. Every number is fitted to the Wave 1 decomposition.'
ans['A1'].font = F_BOLD
rows = [
    (3, 'Regulars: sign-ups a week', round(ans_base)),
    (4, 'Brakes: sign-ups lost a week', round(ans_brake)),
    (5, 'Sale: extra sign-ups in a sale week', round(ans_sale, -1)),
    (6, 'RFL adverts: sign-ups per £1,000', round(ans_rfl, 1)),
    (7, 'RFL adverts: memory %', mem_rfl),
    (8, 'Committed Giving halo: sign-ups per £1,000', round(ans_cog, 1)),
    (9, 'Legacy halo: sign-ups per £1,000', round(ans_leg, 1)),
    (10, 'Brand adverts: sign-ups per £1,000', round(ans_brand, 1)),
    (11, 'Tolerance on a dial (share either side)', 0.25),
    (12, 'Tolerance on memory (points either side)', 10),
    (13, 'Tolerance on a channel dial (share either side)', 0.35),
    (14, 'Brand adverts memory %', mem_brand),
    (15, '(spare)', ''),
]
for r, lab, v in rows:
    ans[f'B{r}'] = lab
    ans[f'C{r}'] = v
for j, ((g, _, _), (m, k)) in enumerate(zip(channel_groups, ch_fits)):
    ans[f'B{16 + j}'] = f'{g}: memory %'
    ans[f'C{16 + j}'] = m
    ans[f'D{16 + j}'] = f'{g}: sign-ups per £1,000'
    ans[f'E{16 + j}'] = round(k, 1)
ans['B22'] = 'Level 6 split in use? (all five channel dials typed)'
ans['C22'] = '=IF(COUNT(' + ','.join(D['ch']) + ')=5,"Yes","No")'
ans['B24'] = 'Pattern-explained targets by level'
for lv in (3, 4, 5, 6):
    ans[f'B{24 + lv}'] = f'Level {lv} target'
    ans[f'C{24 + lv}'] = round(max(0.0, targets[lv] - 0.05), 2)
ans['B33'] = "The answer's own scores (for the honesty note on Level 7)"
ans['C33'] = round(ans_r2, 3)
ans['C34'] = round(ans_mape, 3)
ans['C35'] = round(ans_dw, 2)
ans['B34'] = 'Answer MAPE'
ans['B35'] = 'Answer Durbin-Watson'
ans.column_dimensions['B'].width = 52
ans.column_dimensions['D'].width = 44
ans.sheet_state = 'veryHidden'

# ----------------------------------------------------------------------------------
# 5. Level sheets
# ----------------------------------------------------------------------------------
LEVEL_META = {}


def level_sheet(n, tab_title, sub, look, dials, learned, room, next_text, color, chart_builder, checks,
                extra=None, reveal_lines=None, targets_lv=None, next_row=36):
    ws = wb.create_sheet(f'Level {n}')
    ws.sheet_properties.tabColor = color
    set_widths(ws, {'A': 2, 'B': 46, 'C': 17, 'D': 30, 'E': 2})
    title(ws, tab_title, sub)
    header(ws, 4, '1. Look at the picture')
    r = 5
    for t in look:
        text_row(ws, r, t)
        r += 1
    header(ws, 9, '2. Try: type in the blue cells')
    ws.row_dimensions[10].height = 6
    for j, (lab, unit, kind, default) in enumerate(dials):
        rr = 11 + j
        label(ws, f'B{rr}', lab)
        ws.row_dimensions[rr].height = 34
        if kind == 'memory':
            memory_dropdown(ws, f'C{rr}')
        elif kind == 'number':
            input_cell(ws, f'C{rr}', default, fmt='#,##0.0' if 'per £1,000' in lab else '#,##0')
        hint(ws, f'D{rr}', unit)
    header(ws, 15, '3. Check')
    rr = 16
    for lab, formula, fmt in checks:
        label(ws, f'B{rr}', lab)
        c = ws[f'C{rr}']
        c.value = formula
        c.font = F_BOLD
        c.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        c.border = BOX
        if fmt:
            c.number_format = fmt
        ws.row_dimensions[rr].height = 22
        rr += 1
    verdict_formats(ws, f'C16:C{rr}')
    header(ws, 23, '4. What you just learned')
    r = 24
    for t in learned:
        text_row(ws, r, t)
        r += 1
    header(ws, 28, 'Say this in the room')
    text_row(ws, 29, room, font=F_BOLD, fill=FILL_ACCENT)
    label(ws, 'B31', 'Show the answer? Try a few numbers first.', bold=True)
    yes_no(ws, 'C31')
    if reveal_lines:
        rr = 32
        for lab, formula, fmt in reveal_lines:
            label(ws, f'B{rr}', lab)
            c = ws[f'C{rr}']
            c.value = formula.replace('=IF(C31="Yes",', '=IF(AND(C31="Yes",C11<>""),', 1)
            c.font = F_BOLD
            c.alignment = Alignment(horizontal='center')
            if fmt:
                c.number_format = fmt
            rr += 1
    label(ws, 'D31', 'The answer only appears once you have typed a first number.')
    ws['D31'].font = F_SMALL
    header(ws, next_row, 'Next')
    text_row(ws, next_row + 1, next_text)
    if extra:
        extra(ws)
    chart_builder(ws)
    ws.freeze_panes = 'A4'
    LEVEL_META[n] = ws
    return ws


def dial_verdict(dial_ref, ans_ref, tol_ref=A_ + '$C$11', low='Too low', high='Too high'):
    return (f'=IF({dial_ref}="","Type a number",IF({dial_ref}<{ans_ref}*(1-{tol_ref}),"{low}",'
            f'IF({dial_ref}>{ans_ref}*(1+{tol_ref}),"{high}","About right")))')


def memory_verdict(dial_ref, ans_ref):
    return (f'=IF({dial_ref}="","Pick a memory",IF({dial_ref}<{ans_ref}-{A_}$C$12,"Too short: the effect fades too fast",'
            f'IF({dial_ref}>{ans_ref}+{A_}$C$12,"Too long: the effect lingers too much","About right")))')


def r2_formula(model_rng):
    return f'=MAX(0,1-SUMXMY2({ACT},{model_rng})/DEVSQ({ACT}))'


def status_formula(verdict_cells):
    conds = ','.join(f'{c}="About right"' for c in verdict_cells)
    return f'=IF(AND({conds}),"Level done. Move on.","Keep going")'


# ---- Level 1
def chart_l1(ws):
    model_chart(ws, 'F4', 1, [(C['Answer line L1'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 1'], 'Your model', CYAN, None)],
                'Real Race for Life sign-ups each week, and your model so far')


level_sheet(
    1, 'Level 1: A few hundred people a week sign up whatever we do',
    'The regulars. One number to type. Two minutes.',
    ["This is the real line: how many people signed up for Race for Life each week for nearly three years. The tall spikes are spring. From August to December almost nobody signs up.",
     "Look closely at the quiet months. Even then a few hundred people sign up every week. They already know us. Many of them get our emails.",
     "Your job on this level is to guess how many people sign up in a week when there are no adverts and no sale. That is the first block of any model."],
    [('People who sign up in a week with no adverts and no sale', 'Look at the quiet autumn weeks. A number in the hundreds.', 'number', None)],
    ["Some people would sign up even if we did nothing. The model counts them first, before it gives any credit to adverts.",
     "Economists call this block the baseline. For Race for Life it is small: under a thousand a week. Almost every other sign-up was caused by something we did.",
     "Wave 1 found this block came from two things: people who already consider Cancer Research UK, and our own emails."],
    'Race for Life has a small base of a few hundred sign-ups a week. Nearly every other sign-up was caused by something we did.',
    'Go to Level 2. Some blocks pull sign-ups down, not up.',
    '5BC5F2', chart_l1,
    [('Your number', dial_verdict(D['base'], A_ + '$C$3'), None),
     ('This level', status_formula(['C16']), None)],
    reveal_lines=[('Wave 1 answer: sign-ups a week', f'=IF(C31="Yes",{A_}C3,"")', '#,##0')])

# ---- Level 2
def chart_l2(ws):
    model_chart(ws, 'F4', 2, [(C['Answer line L2'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 2'], 'Your model', CYAN, None)],
                'Your model after two blocks. It dips below zero in the autumn, and that is fine for now')


level_sheet(
    2, 'Level 2: A dearer entry fee and a shut sign-up window cost sign-ups',
    'The brakes. Blocks can pull down as well as push up.',
    ["Some things push sign-ups up. Some pull them down. The price of entry pulls down. So does the sign-up window being shut in the autumn.",
     "A model needs the pulling-down blocks too. Without them it would give the adverts credit for sign-ups that price and timing decided.",
     "Your line will dip below zero in the quiet months after this level. That is fine. The sale and the adverts will lift it back. A model is only judged when every block is in."],
    [('Sign-ups lost in a typical week to price and the shut window', 'About the same size as the regulars, pulling the other way.', 'number', None)],
    ["Drivers can be negative. A dearer entry fee and a shut window both cost sign-ups, and the model measures how many.",
     "Economists call these negative drivers. Price is the classic one. In 2025 entry fees were 12% lower on average, and Wave 1 found that added about 2,500 sign-ups.",
     "Wave 1's price block also does a second job: it tells the model that sign-ups are open in spring and closed in autumn."],
    'Price pulls sign-ups down. The model measures the brake as well as the accelerator, or it would over-credit the adverts.',
    'Go to Level 3. The tall spikes are sale weeks.',
    '5BC5F2', chart_l2,
    [('Your number', dial_verdict(D['brake'], A_ + '$C$4'), None),
     ('This level', status_formula(['C16']), None)],
    reveal_lines=[('Wave 1 answer: sign-ups lost a week', f'=IF(C31="Yes",{A_}C4,"")', '#,##0')])

# ---- Level 3
def chart_l3(ws):
    model_chart(ws, 'F4', 3, [(C['Answer line L3'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 3'], 'Your model', CYAN, None)],
                'Your model after three blocks. The spikes should start to appear')


level_sheet(
    3, 'Level 3: A sale week brings in thousands, and some of them were coming anyway',
    'The sale. The biggest single block, with a twist.',
    ["Look at the tall spikes on the real line. Every one is a discount week. Sales are the biggest single thing that moves sign-ups.",
     "Now look at the week after a sale. It dips. Some people who would have signed up next week signed up early to get the deal.",
     "Type how many extra sign-ups a typical sale week brings. The model already knows which weeks were sale weeks and that the end-of-quarter sales were about twice as strong."],
    [('Extra sign-ups in a typical sale week', 'How tall is a spike above the weeks around it? Thousands.', 'number', None)],
    ["A sale adds sign-ups. It also moves some sign-ups earlier. The model counts the dip after the sale so it does not over-credit the sale.",
     "Economists call the early sign-ups pull-forward. Wave 1 added it after the first results, and it made the media numbers more honest.",
     "Discount windows explained 21% of 2025 sign-ups. That is the biggest block after the adverts."],
    'Sale weeks drove a fifth of sign-ups. The dip after each sale is real, and the model counts it.',
    'Go to Level 4. Now the adverts, and the kettle.',
    '5BC5F2', chart_l3,
    [('Your number', dial_verdict(D['sale'], A_ + '$C$5'), None),
     ('Pattern explained so far', r2_formula(rng('AL')), '0%'),
     ('Target for this level', f'={A_}C27', '0%'),
     ('This level', status_formula(['C16']), None)],
    reveal_lines=[('Wave 1 answer: extra sign-ups in a sale week', f'=IF(C31="Yes",{A_}C5,"")', '#,##0')])

# ---- Level 4
def chart_l4(ws):
    model_chart(ws, 'F4', 4, [(C['Answer line L4'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 4'], 'Your model', CYAN, None)],
                'Your model after four blocks against the real line')
    bar = BarChart()
    bar.type = 'col'
    bar.title = 'Spend on Race for Life adverts each week (grey bars) and the sign-ups your block gives (blue line)'
    bar.style = 2
    bar.width = 24
    bar.height = 9
    bar.legend.position = 'b'
    bar.y_axis.title = 'Spend, £ thousands'
    bar.y_axis.number_format = '#,##0'
    bar.y_axis.majorGridlines = None
    bar.x_axis.number_format = 'mmm yy'
    bar.x_axis.delete = False
    bar.y_axis.delete = False
    bar.x_axis.tickLblSkip = 13
    bar.gapWidth = 30
    ref = Reference(wb[ENG], min_col=C['RFL adverts £k (for chart)'], min_row=R0 - 1, max_row=R1)
    bar.add_data(ref, titles_from_data=True)
    bar.series[0].graphicalProperties.solidFill = 'C6C6C6'
    bar.series[0].graphicalProperties.line.solidFill = 'C6C6C6'
    name_series(bar.series[0], 'Spend on Race for Life adverts, £ thousands')
    bar.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
    line = LineChart()
    ref2 = Reference(wb[ENG], min_col=C['Your block 4 on this level 4 chart'], min_row=R0 - 1, max_row=R1)
    line.add_data(ref2, titles_from_data=True)
    name_series(line.series[0], 'Sign-ups a week from your adverts block')
    line_style(line.series[0], CYAN)
    line.y_axis.axId = 200
    line.y_axis.title = 'Sign-ups a week from adverts'
    line.y_axis.number_format = '#,##0'
    line.y_axis.crosses = 'max'
    line.y_axis.majorGridlines = None
    line.y_axis.delete = False
    bar += line
    ws.add_chart(bar, 'F27')


def extra_l4(ws):
    label(ws, 'B13', 'That is this much sponsorship income for every £1 spent', bold=False)
    ws['C13'] = f'=IF(C11="","",C11*{POUND}/1000)'
    ws['C13'].number_format = '"£"0.00'
    ws['C13'].font = F_BOLD
    ws['C13'].alignment = Alignment(horizontal='center')
    hint(ws, 'D13', 'Wave 1 counts about £90 of sponsorship per sign-up. Their figure: about £1.50 across the three years, £1.66 in 2025 alone.')
    ws.row_dimensions[13].height = 34


level_sheet(
    4, 'Level 4: Adverts keep working for weeks after they stop, like a kettle stays hot',
    'The kettle. Two numbers: how much, and how long the memory lasts.',
    ["The grey bars on the lower picture are what we spent each week on Race for Life adverts: TV, radio, online video, posters, social and the rest.",
     "Sign-ups do not jump only in the week an advert runs. They carry on for weeks afterwards. Think of a kettle. Switch it off and the water stays hot for a while.",
     "Two numbers to type. First, how many sign-ups every £1,000 brings in total, counting all the weeks it keeps working. Second, the memory: how much of last week's effect is still there this week."],
    [('Sign-ups from every £1,000 spent, counting every week it keeps working', 'Somewhere between 5 and 30.', 'number', None),
     ('Memory: how much of last week\'s effect is still there this week', '0% means the advert only works in its own week. 90% means it fades very slowly. Pick from the list.', 'memory', None)],
    ["Adverts work in the week they run and keep working after. The memory number says how fast that fades. Economists call the fading adstock.",
     "Judge a campaign in the week it airs and you miss most of what it earned. Wave 1 found the short-term effect lasts up to six months.",
     "Wave 1 found Race for Life adverts returned £1.66 of sponsorship income for every £1 in 2025, and about half of all sign-ups came from them."],
    'Our adverts keep raising sign-ups for weeks after they stop. Judge a campaign in week one and you miss most of what it earned.',
    'Go to Level 5. Adverts for other causes help Race for Life too.',
    '009CEE', chart_l4,
    [('Sign-ups per £1,000', dial_verdict(D['rfl'], A_ + '$C$6'), None),
     ('Memory', memory_verdict(D['rfl_mem'], A_ + '$C$7'), None),
     ('Pattern explained so far', r2_formula(rng('AM')), '0%'),
     ('Target for this level', f'={A_}C28', '0%'),
     ('This level', status_formula(['C16', 'C17']), None)],
    extra=extra_l4,
    reveal_lines=[('Wave 1 answer: sign-ups per £1,000', f'=IF(C31="Yes",{A_}C6,"")', '#,##0.0'),
                  ('Wave 1 answer: memory', f'=IF(C31="Yes",{A_}C7,"")', '0"%"')])

# ---- Level 5
JAN_OCT_24 = f'SUMIFS({{col}},{YEAR},2024,{MONTH},"<=10")'
JAN_OCT_25 = f'SUMIFS({{col}},{YEAR},2025,{MONTH},"<=10")'


def chart_l5(ws):
    model_chart(ws, 'F4', 5, [(C['Answer line L5'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 5'], 'Your model', CYAN, None)],
                'Your model after five blocks against the real line')
    ch = LineChart()
    base_chart(ch, 'The halo blocks on their own: sign-ups a week from adverts that never mention Race for Life', height=9)
    ref = Reference(wb[ENG], min_col=C['Your halo blocks'], min_row=R0 - 1, max_row=R1)
    ch.add_data(ref, titles_from_data=True)
    name_series(ch.series[0], 'Your three halo blocks together')
    line_style(ch.series[0], '00007E')
    ch.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
    ws.add_chart(ch, 'F27')


def extra_l5(ws):
    header(ws, 39, 'What if: the 2025 cut')
    text_row(ws, 40, "In 2025 we spent 42% less on Committed Giving adverts than in 2024. Your model puts a number on what that cost Race for Life, January to October, against the same months of 2024.")
    label(ws, 'B41', 'Race for Life sign-ups lost, your model')
    ws['C41'] = f'=IF(C11="","",{JAN_OCT_24.format(col=rng("AB"))}-{JAN_OCT_25.format(col=rng("AB"))})'
    ws['C41'].number_format = '#,##0'
    ws['C41'].font = F_BOLD
    ws['C41'].alignment = Alignment(horizontal='center')
    label(ws, 'B42', 'Race for Life sign-ups lost, Wave 1')
    ws['C42'] = 9139
    ws['C42'].number_format = '#,##0'
    ws['C42'].font = F_BOLD
    ws['C42'].alignment = Alignment(horizontal='center')
    hint(ws, 'D42', "AP's waterfall, 2024 to 2025. Second biggest cause of the 7% fall, after the Race for Life cut itself.")
    label(ws, 'B43', 'Now try: cut Committed Giving adverts by a further')
    input_cell(ws, 'C43', 20, fmt='0"%"')
    hint(ws, 'D43', 'Type a percentage of the 2025 spend.')
    label(ws, 'B44', 'Race for Life sign-ups that would go with it, a full year')
    ws['C44'] = f'=IF(C11="","",C11*C43/100*SUMIFS({rng("J")},{YEAR},2025))'
    ws['C44'].number_format = '#,##0'
    ws['C44'].font = F_BOLD
    ws['C44'].alignment = Alignment(horizontal='center')
    hint(ws, 'D44', 'Based on 2025 spend of about £3.2m and your halo number.')
    label(ws, 'B45', 'Sponsorship income that would go with it, £k')
    ws['C45'] = f'=IF(C44="","",C44*{POUND}/1000)'
    ws['C45'].number_format = '"£"#,##0"k"'
    ws['C45'].font = F_BOLD
    ws['C45'].alignment = Alignment(horizontal='center')
    ws.row_dimensions[45].height = 30
    for rr in (41, 42, 43, 44):
        ws.row_dimensions[rr].height = 30


level_sheet(
    5, 'Level 5: Adverts for other causes water the Race for Life field too',
    'The halo. Rain on one field waters the next field.',
    ["Committed Giving adverts ask people to give monthly. Legacy adverts ask people to remember us in their will. Brand adverts just say who we are. None of them mentions Race for Life.",
     "Yet when they ran, Race for Life sign-ups rose. Rain on one field waters the next field too.",
     "Three numbers to type: how many Race for Life sign-ups every £1,000 of each kind of advert brings. They are small per pound. The spend is large, so they add up."],
    [('Committed Giving adverts: Race for Life sign-ups per £1,000', 'Small. Under 10.', 'number', None),
     ('Legacy adverts: Race for Life sign-ups per £1,000', 'Small. Under 10.', 'number', None),
     ('Brand adverts: Race for Life sign-ups per £1,000', 'Smaller still. Under 5.', 'number', None)],
    ["Adverts for one product help another. Economists call this a halo. It is measured from the data, not assumed.",
     "In 2023 the Committed Giving halo was 13% of Race for Life sign-ups. By 2025 it was 5%, because the spend was cut by 42%.",
     "The halo is smaller than it first looked. AP rebuilt it after feedback that it felt too big, and now measure Committed Giving as one block rather than channel by channel."],
    'Committed Giving adverts were quietly filling Race for Life too. When we cut them, Race for Life felt it.',
    'Go to Level 6. Split the Race for Life adverts into channels and meet the catch.',
    '009CEE', chart_l5,
    [('Committed Giving', dial_verdict(D['cog'], A_ + '$C$8'), None),
     ('Legacy', dial_verdict(D['leg'], A_ + '$C$9'), None),
     ('Brand', dial_verdict(D['brand'], A_ + '$C$10'), None),
     ('Pattern explained so far', r2_formula(rng('AN')), '0%'),
     ('Target for this level', f'={A_}C29', '0%'),
     ('This level', status_formula(['C16', 'C17', 'C18']), None)],
    extra=extra_l5, next_row=46,
    reveal_lines=[('Wave 1 answer: Committed Giving per £1,000', f'=IF(C31="Yes",{A_}C8,"")', '#,##0.0'),
                  ('Wave 1 answer: Legacy per £1,000', f'=IF(C31="Yes",{A_}C9,"")', '#,##0.0'),
                  ('Wave 1 answer: Brand per £1,000', f'=IF(C31="Yes",{A_}C10,"")', '#,##0.0')])

# ---- Level 6
def chart_l6(ws):
    model_chart(ws, 'F4', 6, [(C['Answer line L6'], 'Wave 1 answer', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your model after level 6'], 'Your model', CYAN, None)],
                'Your model with the adverts split into five channels')
    ch = LineChart()
    base_chart(ch, 'Weekly spend by channel, £ thousands. See how often they rise and fall together', y_title='Spend, £ thousands', height=9)
    for col, color in zip(['TV £k', 'Online video £k', 'Radio £k', 'Posters £k', 'Digital and rest £k'],
                          ['00007E', CYAN, MAGENTA, 'C75000', '878787']):
        ref = Reference(wb[ENG], min_col=C[col], min_row=R0 - 1, max_row=R1)
        ch.add_data(ref, titles_from_data=True)
        name_series(ch.series[-1], col.replace(' £k', ''))
        line_style(ch.series[-1], color, width_pt=1.25)
    ch.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
    ws.add_chart(ch, 'F27')


def extra_l6(ws):
    # memory shown per channel
    for j in range(5):
        hint(ws, f'D{11 + j}', f'="Memory fixed at "&{A_}C{16 + j}&"%. Sign-ups per £1,000, counting every week."')
    header(ws, 39, 'How sure can we be? Three things the data tells us')
    text_row(ws, 40, "How often the channels move together. A number near 1 means they rise and fall in the same weeks, so the model struggles to tell them apart. Near 0 means they take turns, which is what a model needs.")
    pairs = [('TV and online video', 'TV £k', 'Online video £k'), ('TV and radio', 'TV £k', 'Radio £k'),
             ('Radio and posters', 'Radio £k', 'Posters £k'), ('Online video and digital', 'Online video £k', 'Digital and rest £k')]
    for j, (lab, a, b) in enumerate(pairs):
        rr = 41 + j
        label(ws, f'B{rr}', lab)
        ws[f'C{rr}'] = f'=CORREL({rng(CL[a])},{rng(CL[b])})'
        ws[f'C{rr}'].number_format = '0.00'
        ws[f'C{rr}'].font = F_BOLD
        ws[f'C{rr}'].alignment = Alignment(horizontal='center')
        ws.row_dimensions[rr].height = 20
    text_row(ws, 46, "Race for Life TV stopped after 2024. That gave the model two seasons with TV and one without. Changes like that are what let a model tell channels apart. A channel that never changes cannot be measured.")
    text_row(ws, 47, "How sure AP were, from their own tests. TV, online video, posters, radio and regional as a group: sure (t-stat 3.1, where above 2 means the link is very unlikely to be luck). Social, search, display, audio and YouTube as a group: very sure (4.6). Posters on their own: the weakest channel, returning about 60p per £1 with a wide range either side.")


level_sheet(
    6, 'Level 6: When channels run in the same weeks, the model is sure about the group and less sure about each one',
    'Who sang louder. Five numbers, and the honest catch in every media model.',
    ["In Level 4 you treated all Race for Life adverts as one block. Now split them into five: TV, online video, radio, posters, and digital and social together.",
     "Here is the catch. Most of them run in the same weeks. Two people singing the same note: you hear the sound clearly but cannot tell who is louder.",
     "Your Level 4 answer is the guide. AP do exactly this: measure the group first, then split it using the weeks where the channels differ. The checks here are looser on purpose."],
    [('TV: sign-ups per £1,000', '', 'number', None),
     ('Online video: sign-ups per £1,000', '', 'number', None),
     ('Radio: sign-ups per £1,000', '', 'number', None),
     ('Posters: sign-ups per £1,000', '', 'number', None),
     ('Digital, social and the rest: sign-ups per £1,000', '', 'number', None)],
    ["Channels that air together are hard to tell apart. Economists call this multicollinearity. It is not a fault in the model. It is a fact about the media plan.",
     "AP handle it by measuring the group first, then splitting it within a range. That is why the group number is firm and each channel comes with a range around it.",
     "Once all five dials are typed, the game uses your split instead of your Level 4 block. Your pattern score barely moves. That is the lesson: the split changes the story, not the fit."],
    'The group number is firm. The channel split is softer. Quote the range when a decision is close.',
    'Go to Level 7. Is your model any good?',
    '009CEE', chart_l6,
    [('TV', dial_verdict(D['ch'][0], A_ + '$E$16', A_ + '$C$13'), None),
     ('Online video', dial_verdict(D['ch'][1], A_ + '$E$17', A_ + '$C$13'), None),
     ('Radio', dial_verdict(D['ch'][2], A_ + '$E$18', A_ + '$C$13'), None),
     ('Posters', dial_verdict(D['ch'][3], A_ + '$E$19', A_ + '$C$13'), None),
     ('Digital, social and the rest', dial_verdict(D['ch'][4], A_ + '$E$20', A_ + '$C$13'), None),
     ('Pattern explained with the split', r2_formula(rng('AO')), '0%'),
     ('This level', status_formula(['C16', 'C17', 'C18', 'C19', 'C20']), None)],
    extra=extra_l6, next_row=49,
    reveal_lines=[('Wave 1 answers: TV, video, radio, posters, digital per £1,000',
                   f'=IF(C31="Yes",TEXT({A_}E16,"0.0")&", "&TEXT({A_}E17,"0.0")&", "&TEXT({A_}E18,"0.0")&", "&TEXT({A_}E19,"0.0")&", "&TEXT({A_}E20,"0.0"),"")', None)])

# ---- Level 7: tests
def chart_l7(ws):
    model_chart(ws, 'F4', 7, [(C['Answer line L7'], 'Wave 1 model', GREY, 'dash'), (C['Real sign-ups'], 'Real sign-ups', MAGENTA, None),
                              (C['Your final model'], 'Your model', CYAN, None)],
                'Your finished model against the real line. Tick Show the answer to see Wave 1 as well')
    ch = LineChart()
    base_chart(ch, 'The misses: real sign-ups minus your model. Random scatter around zero is good. Long runs above or below are not', height=9)
    ref = Reference(wb[ENG], min_col=C['Miss (real minus yours)'], min_row=R0 - 1, max_row=R1)
    ch.add_data(ref, titles_from_data=True)
    name_series(ch.series[0], 'Real minus yours')
    line_style(ch.series[0], '00007E', width_pt=1.25)
    ch.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
    ws.add_chart(ch, 'F27')


ws7 = wb.create_sheet('Level 7')
ws7.sheet_properties.tabColor = '00007E'
set_widths(ws7, {'A': 2, 'B': 40, 'C': 13, 'D': 13, 'E': 2})
title(ws7, 'Level 7: A good model explains the pattern, misses at random and predicts weeks it never saw',
      'Is it any good? Four questions, your score beside Wave 1\'s. Nothing to type.')
header(ws7, 4, '1. Look at the picture')
text_row(ws7, 5, "Your finished model sits on top of the real line. Where the two lines touch, your model has explained that week. Where they part, something is missing or too big.", merge_to='D')
text_row(ws7, 6, "The lower picture shows the misses week by week. A good model misses a little every week and at random. A model that misses in long runs has forgotten a block.", merge_to='D')
header(ws7, 8, '2. Check: your model against Wave 1')
hdr = ['The test', 'Yours', 'Wave 1', 'Verdict', 'What it means in plain words']
for j, h in enumerate(hdr):
    c = ws7.cell(row=9, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
ws7.column_dimensions['E'].width = 13
ws7.column_dimensions['F'].width = 60
ws7.row_dimensions[9].height = 22
FINAL_R = rng('AP')
y23_24 = f'SUMPRODUCT(({YEAR}<2025)'
tests = [
    ('Pattern explained (R squared)', r2_formula(FINAL_R), 0.94, '=IF(C10>=0.85,"Good","Keep building")',
     "How much of the weekly up-and-down your model reproduces. AP's guide is above 85%. Wave 1 reached 94%.", '0%'),
    ('Typical weekly miss (MAPE, weeks over 2,000 sign-ups)', f'=IFERROR(AVERAGE({rng("BT")}),"")', 0.28, '=IF(C11="","",IF(C11<=0.35,"Good","Keep building"))',
     "On an average season week, how far your number is from the real one. Quiet weeks are left out, as AP did, because a miss of 50 on a week of 100 looks like 50%. AP's guide is under 15%. Wave 1 sits at 28%, which AP accept because the seasonal peaks are so sharp.", '0%'),
    ('Are the misses random? (Durbin-Watson)', f'=IFERROR(SUM({rng("BS")})/SUM({rng("BR")}),"")', 1.51, '=IF(C12="","",IF(AND(C12>=1,C12<=3),"Good","Runs of misses"))',
     "Near 2 means the misses are random noise. Well below 1 means the model keeps missing in runs, usually a seasonal block it has not got. Wave 1 scored 1.51. The game's own answer scores about 1.0 because its memory curves are simpler than AP's.", '0.00'),
    ('Fit on 2023 to 2024 (the weeks you tuned on)', f'=MAX(0,1-SUMPRODUCT(({YEAR}<2025)*({ACT}-{FINAL_R})^2)/SUMPRODUCT(({YEAR}<2025)*({ACT}-SUMPRODUCT(({YEAR}<2025)*{ACT})/COUNTIF({YEAR},"<2025"))^2))', None, '=IF(C13>=0.85,"Good","Keep building")',
     "The pattern test on the first two years only.", '0%'),
    ('Holdout: fit on 2025 (weeks a good model predicts)', f'=MAX(0,1-SUMPRODUCT(({YEAR}=2025)*({ACT}-{FINAL_R})^2)/SUMPRODUCT(({YEAR}=2025)*({ACT}-SUMPRODUCT(({YEAR}=2025)*{ACT})/COUNTIF({YEAR},2025))^2))', None, '=IF(C14>=0.8,"Good","Copied the past, not the pattern")',
     "A model can copy the past and still be wrong about the future. If the first two years are green and 2025 is red, the model copied the noise, not the pattern. AP test this with holdouts.", '0%'),
    ('Unexplained change, 2024 to 2025 (January to October)', f'=IFERROR((({JAN_OCT_25.format(col=ACT)}-{JAN_OCT_24.format(col=ACT)})-({JAN_OCT_25.format(col=FINAL_R)}-{JAN_OCT_24.format(col=FINAL_R)}))/{JAN_OCT_24.format(col=ACT)},"")', 0.014, '=IF(C15="","",IF(ABS(C15)<=0.05,"Good","Something is missing"))',
     "Sign-ups fell 7% in 2025. This is the part of that fall your model cannot pin on any block. Wave 1 left 1.4% unexplained, well inside AP's tolerance. Every model has some.", '0.0%'),
]
for j, (lab, f, w1, verdict, meaning, fmt) in enumerate(tests):
    rr = 10 + j
    label(ws7, f'B{rr}', lab)
    ws7[f'C{rr}'] = f
    ws7[f'C{rr}'].number_format = fmt
    ws7[f'C{rr}'].font = F_BOLD
    ws7[f'C{rr}'].alignment = Alignment(horizontal='center', vertical='center')
    if w1 is not None:
        ws7[f'D{rr}'] = w1
        ws7[f'D{rr}'].number_format = fmt
    else:
        ws7[f'D{rr}'] = 'not shown'
        ws7[f'D{rr}'].font = F_SMALL
    ws7[f'D{rr}'].alignment = Alignment(horizontal='center', vertical='center')
    ws7[f'E{rr}'] = verdict
    ws7[f'E{rr}'].font = F_BOLD
    ws7[f'E{rr}'].alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    ws7[f'F{rr}'] = meaning
    ws7[f'F{rr}'].font = F_BODY
    ws7[f'F{rr}'].alignment = WRAP
    for colx in 'BCDEF':
        ws7[f'{colx}{rr}'].border = BOX
    ws7.row_dimensions[rr].height = 15 * max(2, math.ceil(len(meaning) / 58)) + 4
ws7.conditional_formatting.add('E10:E15', FormulaRule(formula=['E10="Good"'], fill=PatternFill('solid', fgColor=GREENFILL), font=Font(name='Arial', bold=True, color=GREEN)))
ws7.conditional_formatting.add('E10:E15', FormulaRule(formula=['AND(E10<>"Good",E10<>"")'], fill=PatternFill('solid', fgColor=AMBERFILL), font=Font(name='Arial', bold=True, color=AMBER)))
label(ws7, 'B17', 'This level', bold=True)
ws7['C17'] = '=IF(AND(E10="Good",E13="Good",E14="Good"),"Level done. Move on.","Keep going")'
ws7['C17'].font = F_BOLD
ws7['C17'].border = BOX
ws7.merge_cells('C17:D17')
ws7['C17'].alignment = Alignment(horizontal='center')
verdict_formats(ws7, 'C17:C17')
header(ws7, 19, '3. What you just learned')
text_row(ws7, 20, "Four questions judge any model. Does it explain most of the weekly up-and-down? How far off is a typical week? Are the misses random? Does it predict weeks it was not shown?", merge_to='F')
text_row(ws7, 21, "Economists call these R squared, MAPE, Durbin-Watson and the holdout test. Wave 1 scored 94%, 28%, 1.51 and left 1.4% of the 2025 fall unexplained. You can now read that line and know what it means.", merge_to='F')
text_row(ws7, 22, f"Honesty note. The game's own answer, with every dial set to Wave 1's numbers, scores about {ans_r2:.0%} on pattern explained and {ans_mape:.0%} on weekly miss. It sits a little under Wave 1 because the game uses one simple memory curve per block where AP fit a curve per channel. The blocks and their sizes are Wave 1's.", merge_to='F', font=F_SMALL)
header(ws7, 24, 'Say this in the room')
text_row(ws7, 25, 'A good model explains most of the weekly pattern, misses at random, predicts weeks it never saw and leaves little of the year-on-year change unexplained. Wave 1 does all four.', font=F_BOLD, fill=FILL_ACCENT, merge_to='F')
label(ws7, 'B31', 'Show Wave 1\'s model on the picture?', bold=True)
yes_no(ws7, 'C31')
header(ws7, 36, 'Next')
text_row(ws7, 37, 'Go to Level 8. Where should the next pound go?', merge_to='F')
chart_l7(ws7)
ws7.freeze_panes = 'A4'
# move charts on level 7 to the right of the wider table
for ch in ws7._charts:
    ch.anchor = 'H4' if ch.anchor == 'F4' else 'H27'
LEVEL_META[7] = ws7

# ---- Level 8: the next pound
ws8 = wb.create_sheet('Level 8')
ws8.sheet_properties.tabColor = 'E60078'
set_widths(ws8, {'A': 2, 'B': 30, 'C': 12, 'D': 11, 'E': 13, 'F': 12, 'G': 12, 'H': 13, 'I': 12, 'J': 12, 'K': 12, 'L': 12, 'M': 13, 'N': 2})
title(ws8, 'Level 8: The average pound returned £1.66 but the last pound returned 82p',
      'The next pound. Move money between channels and see what happens. Same total budget.')
header(ws8, 4, '1. Look at the picture')
text_row(ws8, 5, "The third biscuit never tastes like the first. The first £100,000 on a channel finds the keenest people. The next £100,000 has to persuade harder. So each extra pound earns a little less than the one before.", merge_to='L')
text_row(ws8, 6, "This table is the real 2025 Race for Life plan, January to October, with the return per £1 that Wave 1 measured for each channel. Type a new spend for any channel in the blue cells. Keep the total the same. Try to beat the 2025 plan.", merge_to='L')
header(ws8, 8, '2. Try: move the money')
budget_hdr = ['Channel', '2025 spend £k', 'Wave 1 return per £1', 'Your spend £k', 'Lowest allowed', 'Highest allowed',
              'Plan in use £k', 'Sign-ups', 'Income £k', 'Average return per £1', 'Last-pound return per £1', "AP's plan £k"]
for j, h in enumerate(budget_hdr):
    c = ws8.cell(row=9, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
ws8.row_dimensions[9].height = 44
channels_2025 = [
    ('DRTV', 984684, 1.955, 1375338), ('Online video (VOD)', 899101, 1.861, 768640), ('Radio', 834255, 1.753, 666525),
    ('Paid social', 1062828, 1.591, 1077485), ('Regional', 748218, 1.710, 460481), ('Posters (OOH)', 521408, 0.538, 320322),
    ('Digital audio', 458113, 1.708, 657756), ('Display', 441686, 1.740, 552287), ('Paid search', 317693, 1.7, 374699),
    ('YouTube', 120519, 1.850, 150025), ('Telemarketing', 102187, 1.2, 71033), ('Direct mail', 40000, 1.9, 57392),
    ('Door drops', 22184, 2.0, 20892),
]
first, last = 10, 10 + len(channels_2025) - 1
for j, (name, spend, roi, ap_plan) in enumerate(channels_2025):
    rr = first + j
    ws8[f'B{rr}'] = name
    ws8[f'C{rr}'] = round(spend / 1000)
    ws8[f'D{rr}'] = roi
    input_cell(ws8, f'E{rr}', round(spend / 1000), fmt='#,##0')
    ws8[f'F{rr}'] = f'=ROUND(C{rr}*0.5,0)'
    ws8[f'G{rr}'] = f'=ROUND(C{rr}*1.5,0)'
    ws8[f'H{rr}'] = f'=IF($C$26="AP\'s plan",M{rr},MIN(G{rr},MAX(F{rr},IF(E{rr}="",0,E{rr}))))'
    # sign-ups = K * sqrt(spend); K set so that the 2025 plan reproduces Wave 1's return
    ws8[f'I{rr}'] = f'=IF(H{rr}<=0,0,(D{rr}*C{rr}*1000/90)/SQRT(C{rr})*SQRT(H{rr}))'
    ws8[f'J{rr}'] = f'=I{rr}*{POUND}/1000'
    ws8[f'K{rr}'] = f'=IF(H{rr}<=0,0,J{rr}/H{rr})'
    ws8[f'L{rr}'] = f'=0.5*K{rr}'
    ws8[f'M{rr}'] = round(ap_plan / 1000)
    for colx in 'BCDFGHIJKLM':
        ws8[f'{colx}{rr}'].font = F_BODY
        ws8[f'{colx}{rr}'].border = BOX
    for colx in 'CFGHIJM':
        ws8[f'{colx}{rr}'].number_format = '#,##0'
    for colx in 'DKL':
        ws8[f'{colx}{rr}'].number_format = '"£"0.00'
    ws8[f'E{rr}'].border = BOX
tot = last + 1
ws8[f'B{tot}'] = 'Total'
ws8[f'B{tot}'].font = F_BOLD
for colx in 'CEHIJM':
    ws8[f'{colx}{tot}'] = f'=SUM({colx}{first}:{colx}{last})'
    ws8[f'{colx}{tot}'].font = F_BOLD
    ws8[f'{colx}{tot}'].number_format = '#,##0'
    ws8[f'{colx}{tot}'].border = BOX
ws8[f'K{tot}'] = f'=J{tot}/H{tot}'
ws8[f'K{tot}'].number_format = '"£"0.00'
ws8[f'K{tot}'].font = F_BOLD
ws8[f'L{tot}'] = f'=0.5*K{tot}'
ws8[f'L{tot}'].number_format = '"£"0.00'
ws8[f'L{tot}'].font = F_BOLD
text_row(ws8, tot + 1, "Sign-ups use a simple curve: double the spend and you get 41% more sign-ups, not double. Each channel's curve is set so that the 2025 plan gives exactly Wave 1's measured return. AP's real curves differ by channel, so their optimised plan (last column) gained 12,300 sign-ups with them and only about 600 here. The direction is the lesson: money moves towards the channels whose last pound still earns the most.", font=F_SMALL, merge_to='L')
header(ws8, 26, '3. Check')
label(ws8, 'B26', 'Whose plan is on the picture?', bold=True)
dv = DataValidation(type='list', formula1='"Mine,AP\'s plan"', allow_blank=False)
ws8.add_data_validation(dv)
dv.add('C26')
ws8['C26'] = 'Mine'
ws8['C26'].font = F_INPUT
ws8['C26'].fill = FILL_INPUT
ws8['C26'].border = BOX
ws8['C26'].alignment = Alignment(horizontal='center')
hint(ws8, 'D26', "Switch to AP's plan to see where they moved the money: out of radio, posters, regional and video, into DRTV, audio and display.")
ws8.merge_cells('D26:L26')
label(ws8, 'B27', 'Budget check')
ws8['C27'] = f'=IF(ABS(H{tot}-C{tot})<=10,"Same total. Good.",IF(H{tot}>C{tot},"Over budget by £"&TEXT(H{tot}-C{tot},"#,##0")&"k","Under budget by £"&TEXT(C{tot}-H{tot},"#,##0")&"k"))'
ws8.merge_cells('C27:F27')
label(ws8, 'B28', 'Sign-ups from the 2025 plan')
ws8['C28'] = f'=SUMPRODUCT(D{first}:D{last},C{first}:C{last})*1000/90'  # Wave 1 baseline at AP's £90
label(ws8, 'B29', 'Sign-ups from the plan in use')
ws8['C29'] = f'=I{tot}'
label(ws8, 'B30', 'Extra sign-ups you found')
ws8['C30'] = '=C29-C28'
label(ws8, 'B31', 'Extra income, £k')
ws8['C31'] = f'=C30*{POUND}/1000'
label(ws8, 'B32', 'This level')
ws8['C32'] = '=IF(AND(C30>=1000,ABS(H23-C23)<=10),"Level done. Move on.",IF(C30>0,"Keep going: find 1,000 more","Keep going"))'
ws8.merge_cells('C32:F32')
for rr in range(27, 33):
    ws8[f'C{rr}'].font = F_BOLD
    ws8[f'C{rr}'].alignment = Alignment(horizontal='center')
    ws8[f'C{rr}'].number_format = '#,##0'
    ws8.row_dimensions[rr].height = 20
hint(ws8, 'D28', 'Wave 1: 122,189 Race for Life sign-ups from Race for Life media in 2025 (this table covers the thirteen channels with a 2025 return).')
ws8.merge_cells('D28:L28')
hint(ws8, 'D30', "AP's optimiser found 12,300 extra sign-ups (£1.12m) on the same budget with their curves. The most this table can find is about 2,000. Beat 1,000 and you have the idea.")
ws8.merge_cells('D30:L30')
verdict_formats(ws8, 'C32:C32')
ws8.conditional_formatting.add('C27', FormulaRule(formula=['ISNUMBER(SEARCH("Good",C27))'], fill=PatternFill('solid', fgColor=GREENFILL), font=Font(name='Arial', bold=True, color=GREEN)))
ws8.conditional_formatting.add('C27', FormulaRule(formula=['NOT(ISNUMBER(SEARCH("Good",C27)))'], fill=PatternFill('solid', fgColor=AMBERFILL), font=Font(name='Arial', bold=True, color=AMBER)))
header(ws8, 34, '4. What you just learned')
text_row(ws8, 35, "Each extra pound earns a little less than the one before. Economists call this diminishing returns. The return on the very last pound is the marginal return.", merge_to='L')
text_row(ws8, 36, "Wave 1 found the average pound of Race for Life media returned £1.66 in 2025, but the last pound returned 82p. Growth decisions need the last-pound number, not the average. Moving money between channels was worth about £1.1m on the same budget.", merge_to='L')
text_row(ws8, 37, "Posters returned 60p per £1 in the short term, the weakest channel. Over three years, once the slow brand effect is counted, every channel earns a little more. That is Level 9.", merge_to='L')
header(ws8, 39, 'Say this in the room')
text_row(ws8, 40, 'The average pound returned £1.66 but the last pound returned 82p. Moving money between channels is worth about £1.1m before we add a single pound.', font=F_BOLD, fill=FILL_ACCENT, merge_to='L')
header(ws8, 42, 'Next')
text_row(ws8, 43, 'Go to Level 9. The long game: brand and consideration. Then Level 10 for Committed Giving.', merge_to='L')
bar = BarChart()
bar.type = 'col'
bar.title = 'Spend by channel, £ thousands: the 2025 plan against the plan in use'
bar.style = 2
bar.width = 26
bar.height = 10
bar.legend.position = 'b'
bar.y_axis.number_format = '#,##0'
bar.y_axis.majorGridlines = None
bar.x_axis.delete = False
bar.y_axis.delete = False
bar.add_data(Reference(ws8, min_col=3, min_row=9, max_row=last), titles_from_data=True)
bar.add_data(Reference(ws8, min_col=8, min_row=9, max_row=last), titles_from_data=True)
bar.set_categories(Reference(ws8, min_col=2, min_row=first, max_row=last))
bar.series[0].graphicalProperties.solidFill = 'C6C6C6'
bar.series[1].graphicalProperties.solidFill = CYAN
name_series(bar.series[0], '2025 plan')
name_series(bar.series[1], 'Plan in use')
ws8.add_chart(bar, 'B45')
ws8.freeze_panes = 'A4'
LEVEL_META[8] = ws8

# ---- Level 9: Brand, the long game
ws9 = wb.create_sheet('Level 9')
ws9.sheet_properties.tabColor = 'E60078'
set_widths(ws9, {'A': 2, 'B': 46, 'C': 17, 'D': 30, 'E': 2})
title(ws9, 'Level 9: Some of today\'s sign-ups come from adverts people saw years ago',
      'Brand, the long game. Consideration, brand equity and the long-term multiplier.')
header(ws9, 4, '1. Look at the picture')
text_row(ws9, 5, "The line on the picture is brand consideration: the share of people who say they would consider supporting Cancer Research UK. It moves slowly. The bars are everything we spent on media each week, across every cause.")
text_row(ws9, 6, "Wave 1 traced 18.8% of today's consideration back to past media. That stock is our brand equity. It is slipping by about a quarter of a point a year.")
text_row(ws9, 7, "Consideration then turns into sign-ups over the following three years. That is the long-term effect, and it sits on top of the short-term returns you met in Level 8.")
header(ws9, 9, '2. Try: the long-term multiplier')
label(ws9, 'B11', 'Extra long-term media, £ millions a year')
input_cell(ws9, 'C11', 1, fmt='"£"0.0"m"')
hint(ws9, 'D11', 'Wave 1: each extra £1m lifts consideration by about 0.4 points, over one to two years. Only spend above today\'s level counts.')
label(ws9, 'B12', 'Consideration gained, percentage points')
ws9['C12'] = '=IF(C11="","",C11*0.4)'
ws9['C12'].number_format = '0.0'
label(ws9, 'B13', 'Against the yearly drift of')
ws9['C13'] = -0.25
ws9['C13'].number_format = '0.00'
hint(ws9, 'D13', 'Consideration falls about 0.25 points a year without extra investment. Your extra £1m roughly covers a year and a half of drift.')
label(ws9, 'B14', 'Race for Life media: £1 returns in the short term')
ws9['C14'] = 1.66
ws9['C14'].number_format = '"£"0.00'
label(ws9, 'B15', 'Race for Life media: £1 returns over three years')
ws9['C15'] = '=C14*1.21'
ws9['C15'].number_format = '"£"0.00'
hint(ws9, 'D15', 'Wave 1 multiplier 1.21 for Race for Life, 1.04 for Committed Giving. The multipliers come from the brand consideration model.')
for rr in range(11, 16):
    ws9[f'C{rr}'].font = F_BOLD if rr != 11 else F_INPUT
    ws9[f'C{rr}'].alignment = Alignment(horizontal='center')
    ws9.row_dimensions[rr].height = 30
header(ws9, 17, '3. Check: which channels carry the most long-term weight')
text_row(ws9, 18, "Wave 1's 2025 returns per £1 by channel, short term and long term. Posters look weakest in the short term and carry the biggest long-term share. Every channel earns a little more once the slow brand effect is counted.")
lt_hdr = ['Channel', 'Short term £', 'Long term £', 'Three-year total £', 'Multiplier']
for j, h in enumerate(lt_hdr):
    c = ws9.cell(row=19, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
ws9.column_dimensions['E'].width = 18
ws9.column_dimensions['F'].width = 12
lt_rows = [('All Race for Life media', 1.66, 0.35, 1.21), ('Posters (OOH)', 0.54, 0.07, 1.14), ('DRTV', 1.95, 0.15, 1.08),
           ('Online video (VOD)', 1.86, 0.09, 1.05), ('Paid social', 1.59, 0.07, 1.04), ('Regional', 1.71, 0.08, 1.05),
           ('Display', 1.74, 0.11, 1.06), ('Digital audio', 1.71, 0.08, 1.05), ('Radio', 1.75, 0.09, 1.05), ('YouTube', 1.85, 0.09, 1.05),
           ('All Committed Giving media', 0.20, 0.01, 1.04)]
for j, (nm, st, lt, mult) in enumerate(lt_rows):
    rr = 20 + j
    ws9[f'B{rr}'] = nm
    ws9[f'C{rr}'] = st
    ws9[f'D{rr}'] = lt
    ws9[f'E{rr}'] = f'=C{rr}+D{rr}'
    ws9[f'F{rr}'] = mult
    for colx in 'BCDEF':
        ws9[f'{colx}{rr}'].font = F_BOLD if j in (0, 10) else F_BODY
        ws9[f'{colx}{rr}'].border = BOX
        if colx in 'CDE':
            ws9[f'{colx}{rr}'].number_format = '"£"0.00'
        if colx == 'F':
            ws9[f'{colx}{rr}'].number_format = '0.00"x"'
        if colx != 'B':
            ws9[f'{colx}{rr}'].alignment = Alignment(horizontal='center')
text_row(ws9, 31, "Source: Wave 1 results, 2025 long-term ROIs for Race for Life (multipliers measured from the brand consideration model) and Committed Giving. Race for Life's total multiplier is higher than any single channel because it carries over from Race for Life TV, which did not run in 2025.", font=F_SMALL)
header(ws9, 33, '4. What you just learned')
text_row(ws9, 34, "Adverts do two jobs. They raise sign-ups now, and they raise the share of people who would consider us at all. Economists call the second job brand equity, and Wave 1 measures it through consideration.")
text_row(ws9, 35, "Counting the long-term effect takes Race for Life media from £1.66 to £2.01 per pound. The long-term share is understated: the model cannot see Legacy income or existing givers, and Race for Life's base is close to zero, so consideration has little to lift.")
text_row(ws9, 36, "A bigger brand also makes every other advert work harder. That is why the Committed Giving and Legacy halos in Level 5 exist at all.")
header(ws9, 38, 'Say this in the room')
text_row(ws9, 39, 'Nearly a fifth of the people who consider us do so because of adverts they saw years ago. The short-term return is the floor, not the ceiling.', font=F_BOLD, fill=FILL_ACCENT)
header(ws9, 41, 'Next')
text_row(ws9, 42, 'Go to Level 10. Committed Giving, and why 9p per £1 is not the whole story.')
bar = BarChart()
bar.type = 'col'
bar.title = 'All CRUK media spend each week, £ thousands (bars) and brand consideration (line)'
bar.style = 2
bar.width = 24
bar.height = 11
bar.legend.position = 'b'
bar.y_axis.title = 'Spend, £ thousands'
bar.y_axis.number_format = '#,##0'
bar.y_axis.majorGridlines = None
bar.x_axis.number_format = 'mmm yy'
bar.x_axis.delete = False
bar.y_axis.delete = False
bar.x_axis.tickLblSkip = 13
bar.gapWidth = 30
bar.add_data(Reference(wb[ENG], min_col=C['All CRUK media £k'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
bar.series[0].graphicalProperties.solidFill = 'C6C6C6'
bar.series[0].graphicalProperties.line.solidFill = 'C6C6C6'
name_series(bar.series[0], 'All CRUK media, £ thousands a week')
bar.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
line = LineChart()
line.add_data(Reference(wb[ENG], min_col=C['Brand consideration'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
name_series(line.series[0], 'Brand consideration')
line_style(line.series[0], MAGENTA)
line.y_axis.axId = 200
line.y_axis.title = 'Consideration'
line.y_axis.number_format = '0%'
line.y_axis.crosses = 'max'
line.y_axis.majorGridlines = None
line.y_axis.delete = False
line.y_axis.scaling.min = 0.4
line.y_axis.scaling.max = 0.65
bar += line
ws9.add_chart(bar, 'H4')
ch = BarChart()
ch.type = 'col'
ch.grouping = 'stacked'
ch.overlap = 100
ch.title = 'Return per £1 by Race for Life channel in 2025: short term (cyan) and the long-term top-up (magenta)'
ch.style = 2
ch.width = 24
ch.height = 10
ch.legend.position = 'b'
ch.y_axis.number_format = '"£"0.00'
ch.y_axis.majorGridlines = None
ch.x_axis.delete = False
ch.y_axis.delete = False
ch.add_data(Reference(ws9, min_col=3, min_row=19, max_row=29), titles_from_data=True)
ch.add_data(Reference(ws9, min_col=4, min_row=19, max_row=29), titles_from_data=True)
ch.set_categories(Reference(ws9, min_col=2, min_row=20, max_row=29))
ch.series[0].graphicalProperties.solidFill = CYAN
ch.series[1].graphicalProperties.solidFill = MAGENTA
name_series(ch.series[0], 'Short term, up to six months')
name_series(ch.series[1], 'Long term, over three years')
ws9.add_chart(ch, 'H27')
ws9.freeze_panes = 'A4'
LEVEL_META[9] = ws9

# ---- Level 10: Committed Giving
ws10 = wb.create_sheet('Level 10')
ws10.sheet_properties.tabColor = 'E60078'
set_widths(ws10, {'A': 2, 'B': 46, 'C': 17, 'D': 17, 'E': 17, 'F': 30, 'G': 2})
title(ws10, 'Level 10: Committed Giving looks like 9p per £1 because the model only counts the first gift',
      'Committed Giving. What the model saw, what it could not see, and what the 2025 cut cost.')
header(ws10, 4, '1. Look at the picture')
text_row(ws10, 5, "Committed Givers are people who start a monthly gift. Wave 1 modelled how many new ones we won each week, the pink line on the lower picture. The number is small: a few hundred a week, falling through 2025.", merge_to='F')
text_row(ws10, 6, "Marketing drove about 73% of new givers across the three years. In FY25/26 it was 94%, because inflation ate the base and the adverts were nearly all that was left.", merge_to='F')
text_row(ws10, 7, "Yet the return looks tiny: 9p to 20p per £1. Read the whole level before you judge it.", merge_to='F')
header(ws10, 9, '2. Look: three years of new givers')
cog_hdr = ['', 'FY23/24', 'FY24/25', 'FY25/26 to Oct', 'What it means']
for j, h in enumerate(cog_hdr):
    c = ws10.cell(row=10, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
cog_rows = [('New Committed Givers', 12006, 7622, 2993, 'FY25/26 is seven months, April to October.'),
            ('Of which marketing drove', 8568, 5158, 2822, 'Direct adverts plus the halo from Brand, Race for Life, Legacy, Social Challenges and health campaigns.'),
            ('Share driven by marketing', '=C12/C11', '=D12/D11', '=E12/E11', 'Without marketing in FY23/24, 8,568 new givers would not have happened.'),
            ('Committed Giving media spend, £m', 8.4, 4.1, 1.2, 'Direct Committed Giving media only. Nearly quartered.'),
            ('Return per £1 on that spend, short term', 0.18, 0.14, 0.20, 'Income from the new givers in the modelled window only.'),
            ('Share of media-driven givers that came via halo', 0.33, 0.62, 0.72, 'Less direct spend means more reliance on other causes\' adverts.')]
for j, (lab, a, b, c, why) in enumerate(cog_rows):
    rr = 11 + j
    ws10[f'B{rr}'] = lab
    ws10[f'B{rr}'].font = F_BODY
    for colx, v in zip('CDE', (a, b, c)):
        ws10[f'{colx}{rr}'] = v
        ws10[f'{colx}{rr}'].number_format = '0%' if 'Share' in lab else ('"£"0.00' if 'Return' in lab else ('"£"0.0"m"' if '£m' in lab else '#,##0'))
        ws10[f'{colx}{rr}'].alignment = Alignment(horizontal='center')
        ws10[f'{colx}{rr}'].font = F_BODY
    hint(ws10, f'F{rr}', why)
    for colx in 'BCDEF':
        ws10[f'{colx}{rr}'].border = BOX
    ws10.row_dimensions[rr].height = 30
text_row(ws10, 17, "Source: Wave 1 results, Committed Giving business drivers and short-term performance. Halo share is media-driven income from other causes' adverts divided by all media-driven income.", font=F_SMALL, merge_to='F')
header(ws10, 19, '3. Try: what drove the 23% fall in FY25/26, and what a further cut would do')
wf_hdr = ['Cause (FY24/25 to FY25/26, April to October)', 'Change in new givers', 'Share of the fall', 'In plain words']
for j, h in enumerate(wf_hdr):
    c = ws10.cell(row=20, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
ws10.merge_cells('E20:F20')
waterfall = [('Inflation (the economy)', -904, 'Inflation up 4%. The biggest single cause, and outside our control.'),
             ('Committed Giving adverts cut by nearly half', -307, 'Second biggest. The one we chose.'),
             ('Brand adverts cut by nearly 10%', -107, 'Brand adverts feed Committed Giving too.'),
             ('Race for Life adverts cut by 21%', -102, 'The halo runs both ways.'),
             ('Brand consideration slipped 0.5%', -29, 'Small but real.'),
             ('Legacy adverts almost seven times higher', 41, 'A little help back.'),
             ('Less interest in Macmillan', 46, 'A competitor effect the model could see.'),
             ('Health campaign adverts', 91, 'A halo from a campaign that never mentions giving.'),
             ('Social Challenges adverts up 56%', 337, 'The best replacement for the paid social we cut.')]
for j, (lab, v, why) in enumerate(waterfall):
    rr = 21 + j
    ws10[f'B{rr}'] = lab
    ws10[f'C{rr}'] = v
    ws10[f'C{rr}'].number_format = '+#,##0;-#,##0'
    ws10[f'D{rr}'] = f'=C{rr}/3996'
    ws10[f'D{rr}'].number_format = '+0.0%;-0.0%'
    ws10[f'E{rr}'] = why
    ws10[f'E{rr}'].font = F_SMALL
    ws10[f'E{rr}'].alignment = Alignment(wrap_text=True, vertical='center')
    ws10.merge_cells(f'E{rr}:F{rr}')
    for colx in 'BCD':
        ws10[f'{colx}{rr}'].font = F_BODY
        ws10[f'{colx}{rr}'].border = BOX
        if colx != 'B':
            ws10[f'{colx}{rr}'].alignment = Alignment(horizontal='center')
    ws10.row_dimensions[rr].height = 20
ws10['B30'] = 'New givers, FY24/25 to October'
ws10['C30'] = 3996
ws10['B31'] = 'New givers, FY25/26 to October'
ws10['C31'] = '=C30+SUM(C21:C29)'
ws10['D31'] = '=C31/C30-1'
ws10['D31'].number_format = '+0%;-0%'
hint(ws10, 'E31', "Wave 1 reports 3,061, a 23% fall. The causes above sum to it within rounding.")
ws10.merge_cells('E31:F31')
for rr in (30, 31):
    ws10[f'B{rr}'].font = F_BOLD
    ws10[f'C{rr}'].font = F_BOLD
    ws10[f'C{rr}'].number_format = '#,##0'
    ws10[f'C{rr}'].alignment = Alignment(horizontal='center')
    ws10[f'D{rr}'].alignment = Alignment(horizontal='center')
    ws10[f'D{rr}'].font = F_BOLD
label(ws10, 'B33', 'Now try: cut Committed Giving adverts next year by a further')
input_cell(ws10, 'C33', 25, fmt='0"%"')
hint(ws10, 'D33', 'Type a percentage of the FY25/26 spend.')
ws10.merge_cells('D33:F33')
label(ws10, 'B34', 'New givers that would go with it, over the same seven months')
ws10['C34'] = '=IF(C33="","",ROUND(C33/50*307,0))'
hint(ws10, 'D34', 'Wave 1 saw a 50% cut cost 307 new givers. This scales that in a straight line, so treat it as a guide.')
ws10.merge_cells('D34:F34')
label(ws10, 'B35', 'Income lost in the modelled window, £k')
ws10['C35'] = '=IF(C34="","",C34*300/1000)'
hint(ws10, 'D35', 'About £300 per new giver in the window, Wave 1\'s conversion.')
ws10.merge_cells('D35:F35')
label(ws10, 'B36', 'Your assumption: years a new giver keeps giving')
input_cell(ws10, 'C36', 3, fmt='0')
hint(ws10, 'D36', 'Not a Wave 1 finding. Your own guess, to see what it does to the number.')
ws10.merge_cells('D36:F36')
label(ws10, 'B37', 'Income lost if every year counted the same, £k')
ws10['C37'] = '=IF(OR(C35="",C36=""),"",C35*C36)'
label(ws10, 'B38', 'FY25/26 return per £1 on the same assumption')
ws10['C38'] = '=IF(C36="","",E15*C36)'
ws10['C38'].number_format = '"£"0.00'
hint(ws10, 'D38', 'A simple multiplication, to show why AP say the Committed Giving return is understated. Wave 2 is meant to measure it properly.')
ws10.merge_cells('D38:F38')
for rr in range(33, 39):
    ws10.row_dimensions[rr].height = 30
    if rr not in (33, 36):
        ws10[f'C{rr}'].font = F_BOLD
        ws10[f'C{rr}'].alignment = Alignment(horizontal='center')
        if rr in (34,):
            ws10[f'C{rr}'].number_format = '#,##0'
        if rr in (35, 37):
            ws10[f'C{rr}'].number_format = '"£"#,##0"k"'
header(ws10, 40, '4. What you just learned')
text_row(ws10, 41, "The model counts only brand new givers, at roughly £300 of income each in the modelled window. It does not count the years of giving that follow, and it does not count the givers that adverts stop from leaving. That is why 9p is a floor.", merge_to='F')
text_row(ws10, 42, "Cutting Committed Giving media by nearly half cost about 23% of new givers in FY25/26, and inflation cost about the same again. Two separate causes, one falling line. Economists call the first a controllable driver and the second a base driver.", merge_to='F')
text_row(ws10, 43, "Less direct spend means more reliance on other causes' adverts. In FY25/26 nearly three quarters of media-driven givers came via the halo, mostly from Brand and Social Challenges.", merge_to='F')
header(ws10, 45, 'Say this in the room')
text_row(ws10, 46, 'Committed Giving media shows 9p per pound because the model only counts the first gift of brand new givers. The years of giving that follow, and the churn it prevents, are not in the number yet.', font=F_BOLD, fill=FILL_ACCENT, merge_to='F')
header(ws10, 48, 'Next')
text_row(ws10, 49, 'You have built the whole model. Go to Words for the ladder of every term, or back to Start here to see your progress.', merge_to='F')
ch = BarChart()
ch.type = 'bar'
ch.title = 'Why new givers fell 23% in FY25/26: change in new givers by cause'
ch.style = 2
ch.width = 24
ch.height = 10
ch.legend = None
ch.y_axis.number_format = '#,##0'
ch.y_axis.majorGridlines = None
ch.x_axis.delete = False
ch.y_axis.delete = False
ch.add_data(Reference(ws10, min_col=3, min_row=20, max_row=29), titles_from_data=True)
ch.set_categories(Reference(ws10, min_col=2, min_row=21, max_row=29))
ch.series[0].graphicalProperties.solidFill = MAGENTA
ch.series[0].invertIfNegative = False
ch.x_axis.tickLblPos = 'low'
ws10.add_chart(ch, 'H4')
ch2 = LineChart()
base_chart(ch2, 'New Committed Givers each week, real figures', y_title='New givers a week', height=9)
ch2.add_data(Reference(wb[ENG], min_col=C['Real new Committed Givers'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
name_series(ch2.series[0], 'New Committed Givers a week')
line_style(ch2.series[0], MAGENTA, width_pt=1.25)
ch2.set_categories(Reference(wb[ENG], min_col=1, min_row=R0, max_row=R1))
ws10.add_chart(ch2, 'H26')
ws10.freeze_panes = 'A4'
LEVEL_META[10] = ws10

# ----------------------------------------------------------------------------------
# 6. Words (glossary) with four pictures
# ----------------------------------------------------------------------------------
wsg = wb.create_sheet('Words')
wsg.sheet_properties.tabColor = '00007E'
set_widths(wsg, {'A': 2, 'B': 24, 'C': 44, 'D': 44, 'E': 44, 'F': 12, 'G': 12, 'H': 12})
title(wsg, 'Words: every term, climbed one rung at a time', 'The everyday picture first, then what it meant in Wave 1, then the sentence you can say out loud. The word itself comes last.')
for j, h in enumerate(['The word', 'The everyday picture', 'What it meant in Wave 1', 'Say this in the room']):
    c = wsg.cell(row=4, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
wsg.row_dimensions[4].height = 22
glossary = [
    ('Econometrics (Level 1)', 'A detective with three years of weekly till receipts, working out which adverts raised the money. The detective infers. It never witnesses.',
     'Analytic Partners read 148 weeks of Race for Life sign-ups against everything that happened each week, and measured what each thing was worth.',
     'Econometrics is a detective reading three years of weekly receipts to work out which pounds raised the sign-ups.'),
    ('Baseline (Level 1)', 'The customers who walk into a shop because it is there and open. The sign-ups we would still get with every advert switched off.',
     'Small for Race for Life: a few hundred a week, from people who already consider us and from our own emails. AP say the true base is close to zero.',
     'The base is what we would get with the adverts off. Everything above it is what we caused.'),
    ('Driver (Levels 1 to 6)', 'Anything that can move sign-ups up or down. An advert, a sale, the price, the time of year, even inflation.',
     'AP tested 18 marketing drivers and six non-marketing ones and kept those with a reliable link. Weather, PR and website outages had none.',
     'A driver is anything the model found that moves sign-ups.'),
    ('Negative driver (Level 2)', 'A brake. A dearer ticket means fewer people buy it.',
     'Price and the closed autumn window pull sign-ups down. A 12% lower entry fee in 2025 added about 2,500 sign-ups.',
     'Price pulls sign-ups down, and the model measures the brake as well as the accelerator.'),
    ('Response (Levels 1 to 6)', 'The exchange rate. How many sign-ups one unit of a driver buys: one sale week, or £1,000 of adverts.',
     'This is the number you type on each level. AP estimate it from the weekly data, with a range either side.',
     'The response is the exchange rate between a driver and sign-ups.'),
    ('Pull-forward (Level 3)', 'A January sale. People who were going to buy in February buy in January instead, so February is quiet.',
     'Fewer sign-ups in the weeks after a discount window closed. AP added this after the first results so that sales were not over-credited.',
     'Discounts bring sign-ups forward as well as adding new ones. The dip afterwards is real.'),
    ('Adstock (Level 4)', 'A kettle. Switch it off and the water stays hot for a while. Media heat fades over weeks, not minutes.',
     'Each block has a memory. The short-term effect of Race for Life adverts lasts up to six months.',
     'Our adverts keep raising sign-ups for weeks after they stop, and the model measures how fast that fades.'),
    ('Incremental (Level 4)', 'The part of the spike above the base. Only the sign-ups the advert caused, not the ones that happened while it was on.',
     'About half of Race for Life sign-ups were incremental to Race for Life media in 2025.',
     'Incremental means the sign-ups the advert caused, not the sign-ups that happened while it was on.'),
    ('ROI, return per £1 (Levels 4 and 8)', 'Income earned per pound of adverts. £1.66 back for every £1 in.',
     'Sponsorship income only, about £90 per sign-up. Not entry fees, not Gift Aid, and only new sign-ups.',
     'Every £1 of Race for Life media returned £1.66 of sponsorship income in 2025.'),
    ('Halo (Level 5)', 'Rain on one field waters the next field too. Adverts for one product raise sign-ups for another.',
     'Committed Giving media was 13% of Race for Life sign-ups in 2023 and 5% in 2025. Legacy media added 6 to 8%. Measured, not assumed.',
     'Committed Giving adverts were quietly filling Race for Life too. That is the halo.'),
    ('Multicollinearity (Level 6)', 'Two people singing the same note. You hear the sound but cannot tell who is louder.',
     'Most Race for Life channels air in the same weeks. AP measure the group first, then split it using the weeks where channels differ.',
     'When channels air together the model is sure about the group and less sure about each channel.'),
    ('t-stat and p-value (Level 6)', 'How confident we are that a link is real rather than luck. Above 2 means less than a 5% chance it is luck.',
     'Discounts scored 11.5. The Committed Giving halo scored 4.2. Brand consideration scored 0.96, which is why AP flag it as uncertain.',
     'A t-stat above two means the link is very unlikely to be luck.'),
    ('Confidence interval (Level 6)', 'The range the true answer probably sits in. A forecast says 15 to 19 degrees, not 17.',
     'The Committed Giving halo is 0.002 to 0.004 sign-ups per pound: 2 to 4 per £1,000. Quote the range when a decision is close.',
     'Every result is a range, and the range is the honest answer.'),
    ('R squared (Level 7)', 'How much of the weekly up-and-down the model reproduces. 100% would be a perfect copy.',
     'Wave 1 explains 94% of the pattern. AP\'s guide is above 85%.',
     'The model explains 94% of the weekly pattern in sign-ups.'),
    ('MAPE (Level 7)', 'On an average week, how far the model is from the real number, as a percentage.',
     '28% on weeks with more than 2,000 sign-ups. Quiet weeks are left out because tiny weeks make percentages jump around.',
     'On an average season week the model is within about a quarter of the real number.'),
    ('Durbin-Watson (Level 7)', 'A check that the misses are random scatter, not a pattern the model failed to spot. Near 2 is good.',
     'Wave 1 scored 1.51, acceptable. It tells the analysts whether something seasonal is still missing.',
     'It is a housekeeping check that the errors are random.'),
    ('Holdout test (Level 7)', 'Hide the last few months, fit the model on the rest, then see if it predicts what you hid.',
     'The strongest proof a model works. AP use holdouts, bootstrapping and business sense checks together.',
     'A good model predicts weeks it was never shown.'),
    ('Unexplained (Level 7)', 'The part of a year-on-year change no driver accounts for. Every model has some.',
     '1.4% of sign-ups in the 2024 to 2025 comparison, well inside AP\'s tolerance.',
     'The model left 1.4% unexplained, which is small.'),
    ('Diminishing returns (Level 8)', 'The third biscuit never tastes like the first. It is a curve flattening, not a cliff.',
     'AP build a response curve per channel. The optimiser uses them to find the best spread of the same money.',
     'Each channel has a curve, and the next pound should go where the curve is still steep.'),
    ('Marginal ROI (Level 8)', 'The return on the very last pound, not the average pound.',
     'In 2025 the last pound of Race for Life media returned 82p while the average pound returned £1.66.',
     'The average pound returned £1.66 but the last pound returned 82p, so adding budget needs care.'),
    ('Brand equity (Level 9)', 'Reputation. It keeps working after you stop paying for it, and it drains slowly if you stop.',
     '18.8% of today\'s consideration traces back to past media. It slips 0.25 points a year. An extra £1m lifts it about 0.4 points.',
     'Nearly a fifth of the people who consider us do so because of adverts they saw years ago.'),
    ('Base driver (Level 10)', 'The weather for a farmer. It moves the harvest and nobody on the farm controls it.',
     'Inflation cost Committed Giving about 900 new givers in FY25/26, as much as the advertising cut. Price, seasonality and competitors are base drivers too.',
     'Inflation and the advertising cut each cost about a quarter of new givers. One we chose, one we did not.'),
    ('Long-term multiplier (Level 9)', 'The slow second job of an advert: making more people consider us at all, which turns into sign-ups over three years.',
     '1.21 for Race for Life media (£1.66 becomes £2.01) and 1.04 for Committed Giving. Posters and TV carry the biggest share.',
     'Counting the long-term effect takes Race for Life media from £1.66 to £2.01 per pound.'),
]
r = 5
for word, pic, w1, room in glossary:
    for colx, txt in zip('BCDE', (word, pic, w1, room)):
        c = wsg[f'{colx}{r}']
        c.value = txt
        c.font = F_BOLD if colx == 'B' else F_BODY
        c.alignment = WRAP
        c.border = BOX
    longest = max(len(pic), len(w1), len(room))
    wsg.row_dimensions[r].height = 15 * max(2, math.ceil(longest / 44)) + 6
    r += 1
r += 1
header(wsg, r, 'Four pictures')
text_row(wsg, r + 1, 'Each picture is one idea. The data under them is illustrative, built to show the shape, not Wave 1 figures.', font=F_SMALL, merge_to='E')
pic_row = r + 3
# picture data
d0 = pic_row + 40
wsg[f'B{d0}'] = 'Data for the pictures (illustrative)'
wsg[f'B{d0}'].font = F_SMALL
# 1 kettle: TV on for weeks 1-3 then off, effect with 60% memory
wsg[f'B{d0 + 1}'] = 'Week'
wsg[f'C{d0 + 1}'] = 'Advert on air'
wsg[f'D{d0 + 1}'] = 'Sign-ups it causes (60% memory)'
eff = 0
for w in range(1, 17):
    on = 100 if w <= 3 else 0
    eff = on + 0.6 * eff
    wsg[f'B{d0 + 1 + w}'] = w
    wsg[f'C{d0 + 1 + w}'] = on
    wsg[f'D{d0 + 1 + w}'] = round(eff)
# 2 biscuit: spend vs income
wsg[f'F{d0 + 1}'] = 'Spend £k'
wsg[f'G{d0 + 1}'] = 'Income £k'
wsg[f'H{d0 + 1}'] = 'Income if every pound earned like the first'
for i2, s2 in enumerate(range(0, 1100, 100)):
    wsg[f'F{d0 + 2 + i2}'] = s2
    wsg[f'G{d0 + 2 + i2}'] = round(1.66 * 1000 ** 0.5 * s2 ** 0.5)
    wsg[f'H{d0 + 2 + i2}'] = round(3.3 * s2)
# 3 detective: AP's slide 8 toy
wsg[f'J{d0 + 1}'] = 'Week'
for j, nm in enumerate(['Base', 'Economy', 'TV', 'Radio', 'Sale']):
    wsg.cell(row=d0 + 1, column=11 + j, value=nm)
toy = {
    'Base': [1] * 22,
    'Economy': [0, 0, 0] + [1] * 19,
    'TV': [0] * 6 + [1, 1, 1, 0, 0, 1, 1, 1, 0, 0, 1, 1, 1, 0, 0, 0],
    'Radio': [0] * 11 + [1, 1, 1, 0, 0, 1, 1, 1, 0, 0, 0],
    'Sale': [0] * 16 + [1, 1, 1, 0, 0, 0],
}
for w in range(22):
    wsg[f'J{d0 + 2 + w}'] = w + 1
    for j, nm in enumerate(['Base', 'Economy', 'TV', 'Radio', 'Sale']):
        wsg.cell(row=d0 + 2 + w, column=11 + j, value=toy[nm][w])
# 4 halo: two fields (CG spend index vs RFL halo sign-ups by year, Wave 1 real)
wsg[f'Q{d0 + 1}'] = 'Year'
wsg[f'R{d0 + 1}'] = 'Committed Giving media £m'
wsg[f'S{d0 + 1}'] = 'Race for Life sign-ups from its halo'
for j, (y, a, b) in enumerate([('2023', 8.3, 30778), ('2024', 6.6, 23325), ('2025 to Oct', 3.2, 12178)]):
    wsg[f'Q{d0 + 2 + j}'] = y
    wsg[f'R{d0 + 2 + j}'] = a
    wsg[f'S{d0 + 2 + j}'] = b

ch = BarChart()
ch.type = 'col'
ch.title = 'The kettle: the advert runs for three weeks, the sign-ups carry on for ten'
ch.style = 2
ch.width = 16
ch.height = 8
ch.legend.position = 'b'
ch.y_axis.majorGridlines = None
ch.x_axis.delete = False
ch.y_axis.delete = False
ch.add_data(Reference(wsg, min_col=3, min_row=d0 + 1, max_row=d0 + 17), titles_from_data=True)
ch.series[0].graphicalProperties.solidFill = 'C6C6C6'
ln = LineChart()
ln.add_data(Reference(wsg, min_col=4, min_row=d0 + 1, max_row=d0 + 17), titles_from_data=True)
line_style(ln.series[0], MAGENTA)
ch += ln
ch.set_categories(Reference(wsg, min_col=2, min_row=d0 + 2, max_row=d0 + 17))
wsg.add_chart(ch, f'B{pic_row}')

ch2 = LineChart()
ch2.title = 'The third biscuit: income rises fast at first, then flattens'
ch2.style = 2
ch2.width = 16
ch2.height = 8
ch2.legend.position = 'b'
ch2.y_axis.majorGridlines = None
ch2.x_axis.delete = False
ch2.y_axis.delete = False
ch2.y_axis.title = 'Income £k'
ch2.x_axis.title = 'Spend £k'
ch2.add_data(Reference(wsg, min_col=7, min_row=d0 + 1, max_row=d0 + 12), titles_from_data=True)
ch2.add_data(Reference(wsg, min_col=8, min_row=d0 + 1, max_row=d0 + 12), titles_from_data=True)
line_style(ch2.series[0], MAGENTA, smooth=True)
line_style(ch2.series[1], 'C6C6C6', dash='dash')
ch2.set_categories(Reference(wsg, min_col=6, min_row=d0 + 2, max_row=d0 + 12))
wsg.add_chart(ch2, f'D{pic_row}')

ch3 = BarChart()
ch3.type = 'col'
ch3.grouping = 'stacked'
ch3.overlap = 100
ch3.title = 'The detective: TV ran in all three peaks, radio in two, the sale in one. The differences give each away'
ch3.style = 2
ch3.width = 16
ch3.height = 8
ch3.legend.position = 'b'
ch3.y_axis.majorGridlines = None
ch3.x_axis.delete = False
ch3.y_axis.delete = False
ch3.add_data(Reference(wsg, min_col=11, max_col=15, min_row=d0 + 1, max_row=d0 + 23), titles_from_data=True)
for s, colr in zip(ch3.series, ['C6C6C6', '878787', '00007E', CYAN, MAGENTA]):
    s.graphicalProperties.solidFill = colr
ch3.set_categories(Reference(wsg, min_col=10, min_row=d0 + 2, max_row=d0 + 23))
wsg.add_chart(ch3, f'B{pic_row + 17}')

ch4 = BarChart()
ch4.type = 'col'
ch4.title = 'Rain on the next field: less Committed Giving media, fewer Race for Life sign-ups from its halo'
ch4.style = 2
ch4.width = 16
ch4.height = 8
ch4.legend.position = 'b'
ch4.y_axis.majorGridlines = None
ch4.x_axis.delete = False
ch4.y_axis.delete = False
ch4.y_axis.title = 'Committed Giving media £m'
ch4.add_data(Reference(wsg, min_col=18, min_row=d0 + 1, max_row=d0 + 4), titles_from_data=True)
ch4.series[0].graphicalProperties.solidFill = 'C6C6C6'
ln4 = LineChart()
ln4.add_data(Reference(wsg, min_col=19, min_row=d0 + 1, max_row=d0 + 4), titles_from_data=True)
line_style(ln4.series[0], MAGENTA)
ln4.series[0].marker.symbol = 'circle'
ln4.y_axis.axId = 200
ln4.y_axis.crosses = 'max'
ln4.y_axis.title = 'Race for Life sign-ups'
ln4.y_axis.number_format = '#,##0'
ln4.y_axis.majorGridlines = None
ln4.y_axis.delete = False
ch4 += ln4
ch4.set_categories(Reference(wsg, min_col=17, min_row=d0 + 2, max_row=d0 + 4))
wsg.add_chart(ch4, f'D{pic_row + 17}')
wsg.freeze_panes = 'A5'

# ----------------------------------------------------------------------------------
# 7. Start here
# ----------------------------------------------------------------------------------
ws0 = wb.create_sheet('Start here', 0)
ws0.sheet_properties.tabColor = MAGENTA
set_widths(ws0, {'A': 2, 'B': 22, 'C': 46, 'D': 30, 'E': 20, 'F': 2})
title(ws0, 'The Econometrics Game: find out which pounds raised the sign-ups',
      'Version 2, September 2026. CRUK confidential, internal use only. No macros: type in the blue cells and everything updates.')
header(ws0, 4, 'The idea, in one picture')
text_row(ws0, 5, "Imagine a detective with three years of till receipts. Every week the receipts show how many people signed up for Race for Life.", merge_to='E')
text_row(ws0, 6, "The detective also knows what happened each week. Which adverts ran. Whether there was a sale. Whether sign-ups were even open.", merge_to='E')
text_row(ws0, 7, "The detective's job is to work out which of those things caused the sign-ups. That is all econometrics is. Analytic Partners did it for us in 2026 and called it Wave 1. In this game you do it yourself, one block at a time.", merge_to='E')
header(ws0, 9, 'How to play')
text_row(ws0, 10, "1. Type only in the blue cells. 2. Each level adds one block to your model and asks for one or two numbers. 3. A check tells you when a number is about right. Green means move on. 4. Stuck? Set Show the answer to Yes.", merge_to='E')
header(ws0, 12, 'Two routes')
for j, (lab, txt) in enumerate([
        ('Five minutes', 'Level 1, then Level 4, then the What if box on Level 5, then Level 8. You will leave knowing the base, the kettle, the halo and the last pound.'),
        ('Thirty minutes', 'Levels 1 to 10 in order. Each one builds on the last. Finish with Words, which names every term you have already understood.')]):
    label(ws0, f'B{13 + j}', lab, bold=True)
    text_row(ws0, 13 + j, txt, merge_to=None)
    ws0[f'C{13 + j}'].value = txt
    ws0[f'C{13 + j}'].font = F_BODY
    ws0[f'C{13 + j}'].alignment = WRAP
    ws0[f'B{13 + j}'].value = lab
    ws0.merge_cells(f'C{13 + j}:E{13 + j}')
    ws0.row_dimensions[13 + j].height = 34
header(ws0, 16, 'Your progress')
for j, h in enumerate(['Level', 'What you build', 'The word you learn', 'Status']):
    c = ws0.cell(row=17, column=2 + j, value=h)
    c.font = F_WHITE
    c.fill = FILL_HEAD
    c.alignment = CENTER
    c.border = BOX
progress = [
    (1, 'The regulars', 'Baseline', "='Level 1'!C17"),
    (2, 'The brakes', 'Negative driver', "='Level 2'!C17"),
    (3, 'The sale', 'Pull-forward', "='Level 3'!C19"),
    (4, 'The kettle: Race for Life adverts', 'Adstock', "='Level 4'!C20"),
    (5, 'Rain on the next field: halo', 'Halo', "='Level 5'!C21"),
    (6, 'Who sang louder: the channel split', 'Multicollinearity', "='Level 6'!C22"),
    (7, 'Is it any good: four tests', 'R squared, MAPE, holdout', "='Level 7'!C17"),
    (8, 'The next pound: the budget game', 'Marginal return', "='Level 8'!C32"),
    (9, 'The long game: brand and consideration', 'Brand equity', 'Read and try'),
    (10, 'Committed Giving: why 9p is a floor', 'Base driver', 'Read and try'),
]
for j, (lv, what, word, f) in enumerate(progress):
    rr = 18 + j
    ws0[f'B{rr}'] = f'Level {lv}'
    ws0[f'B{rr}'].hyperlink = f"#'Level {lv}'!A1"
    ws0[f'B{rr}'].font = Font(name='Arial', size=11, bold=True, color=STEEL, underline='single')
    ws0[f'C{rr}'] = what
    ws0[f'D{rr}'] = word
    ws0[f'E{rr}'] = f
    for colx in 'BCDE':
        ws0[f'{colx}{rr}'].border = BOX
        if colx != 'B':
            ws0[f'{colx}{rr}'].font = F_BODY
    ws0[f'E{rr}'].font = F_BOLD
    ws0[f'E{rr}'].alignment = Alignment(horizontal='center', wrap_text=True)
    ws0.row_dimensions[rr].height = 20
verdict_formats(ws0, 'E18:E27')
label(ws0, 'B29', 'Levels done', bold=True)
ws0['C29'] = '=COUNTIF(E18:E27,"Level done*")&" of 8 scored levels"'
ws0['C29'].font = F_BIG
header(ws0, 31, 'Say this in the room')
text_row(ws0, 32, 'Econometrics is a detective reading three years of weekly receipts to work out which pounds raised the sign-ups. It can only see what changed.', font=F_BOLD, fill=FILL_ACCENT, merge_to='E')
header(ws0, 34, 'Where the numbers come from')
text_row(ws0, 35, "Analytic Partners, CRUK Wave 1 Econometrics Results (July 2026) and Wave 1 Methodology and Model Documentation (17 July 2026). The weekly figures are AP's own, read from the charts in the results deck: real sign-ups, spend by channel and product, and AP's weekly decomposition of sign-ups by driver, January 2023 to October 2025. The game's answer key is fitted to that decomposition. Sponsorship income is counted at about £90 per sign-up, as AP do.", merge_to='E', font=F_SMALL)
text_row(ws0, 36, "Built on the Analytic Partners training game. Under the bonnet shows every weekly number and every formula. Words explains every term. Nothing is hidden except the answer key, so the game stays a game.", merge_to='E', font=F_SMALL)
label(ws0, 'B38', 'Sponsorship income per Race for Life sign-up, £', bold=True)
input_cell(ws0, 'C38', 90, fmt='"£"0')
hint(ws0, 'D38', "AP's Wave 1 conversion, about £90 of sponsorship per sign-up. Change it and every pound figure in the game moves with it. Sign-ups never change.")
ws0.merge_cells('D38:E38')
ws0.row_dimensions[38].height = 34
ws0.freeze_panes = 'A4'

# ----------------------------------------------------------------------------------
# 8. Room notes (facilitator), hidden
# ----------------------------------------------------------------------------------
wsr = wb.create_sheet('Room notes')
wsr.sheet_properties.tabColor = GREY
set_widths(wsr, {'A': 2, 'B': 26, 'C': 90})
title(wsr, 'Room notes: running the game with a team in 45 minutes', 'For the person leading the session. Unhide this tab from the sheet list. The answers are here, so do not share the screen while it is open.')
plan = [
    ('0 to 5 minutes', "Open Start here. Read the detective paragraph aloud and stop before the last sentence. Ask the room: what do you think caused the spikes? Let three people answer. Do not correct anyone yet."),
    ('5 to 15 minutes', "Levels 1 to 3 together on the big screen, one person typing. Ask for a show of hands before each number: higher or lower than a thousand? The base surprises people. Let it."),
    ('15 to 25 minutes', "Level 4. Ask the room to guess the memory before you touch it. Most people say 20%. Show what 80% does to the line. Name adstock only after the line fits."),
    ('25 to 32 minutes', "Level 5. Type the halo numbers, then read the What if box. Say the 9,139 out loud and stop. Let the room say what it means for next year's budget."),
    ('32 to 38 minutes', "Level 6 in pairs on their own laptops. The point is the catch, not the numbers. Ask: which channel are you least sure of, and why?"),
    ('38 to 45 minutes', "Level 8 as a race: who can find 1,000 extra sign-ups on the same budget first? Close on the average pound versus the last pound. Send everyone to Levels 9 and 10 and Words for later."),
]
r = 4
header(wsr, r, 'Timings')
for lab, txt in plan:
    r += 1
    wsr[f'B{r}'] = lab
    wsr[f'B{r}'].font = F_BOLD
    wsr[f'B{r}'].alignment = Alignment(vertical='top')
    wsr[f'C{r}'] = txt
    wsr[f'C{r}'].font = F_BODY
    wsr[f'C{r}'].alignment = WRAP
    wsr.row_dimensions[r].height = 15 * max(2, math.ceil(len(txt) / 90)) + 4
r += 2
header(wsr, r, 'The answer key')
answers_list = [
    ('Level 1 regulars', f'{ans_base:,.0f} sign-ups a week. About right is 25% either side.'),
    ('Level 2 brakes', f'{ans_brake:,.0f} sign-ups lost a week.'),
    ('Level 3 sale', f'{ans_sale:,.0f} extra sign-ups in a typical sale week. End-of-quarter sales are stronger; the shape knows.'),
    ('Level 4 adverts', f'{ans_rfl:.1f} sign-ups per £1,000 (£{ans_rfl * 0.09:.2f} per £1), memory {mem_rfl}%.'),
    ('Level 5 halo', f'Committed Giving {ans_cog:.1f}, Legacy {ans_leg:.1f}, Brand {ans_brand:.1f} sign-ups per £1,000.'),
    ('Level 6 split', ', '.join(f'{g.split(" (")[0]} {k:.1f}' for (g, _, _), (m, k) in zip(channel_groups, ch_fits)) + ' per £1,000. About right is 35% either side.'),
    ('Level 8 budget', "Any plan that keeps the total and finds 1,000 extra sign-ups. The best possible here is about 2,000: cut posters and telemarketing, add DRTV and online video. Switching to AP's plan shows the direction they moved: into DRTV, audio and display, out of posters, regional, radio and video."),
]
for lab, txt in answers_list:
    r += 1
    wsr[f'B{r}'] = lab
    wsr[f'B{r}'].font = F_BOLD
    wsr[f'C{r}'] = txt
    wsr[f'C{r}'].font = F_BODY
    wsr[f'C{r}'].alignment = WRAP
    wsr.row_dimensions[r].height = 15 * max(1, math.ceil(len(txt) / 90)) + 4
r += 2
header(wsr, r, 'Questions the room will ask, and the plain answer')
traps = [
    ("Isn't this just correlation?", "Partly, and AP say so. The model infers, it never witnesses. Three things make it more than a coincidence: it controls for everything else at once, it is tested on weeks it never saw, and channels that stopped or started (TV stopped after 2024) let it see what happens without them."),
    ("Why is the base so small? Surely people would sign up anyway.", "In the quiet months a few hundred do. In season, nearly every sign-up lines up with a sale or an advert. AP say the true base is close to zero, and that Race for Life is an earned product."),
    ("The halo felt too big.", "It did, and AP rebuilt it. They now measure Committed Giving media as one block rather than by channel, and the halo fell. It is still 5% of 2025 sign-ups and 9,139 of the 2025 decline."),
    ("Why are Committed Giving returns so low?", "The model counts only brand new givers, at about £300 each in the window. Not the years of giving that follow, not the churn that adverts prevent. Level 10 shows the arithmetic. Wave 2 is meant to fix that."),
    ("Can we trust the channel split?", "Trust the group, quote the range on the channel. That is AP's own advice and the whole of Level 6."),
    ("Why not just spend more if £1 returns £1.66?", "Because the last pound returned 82p. The average pound looks great and the next pound does not. Move money first, then add."),
]
for q, a in traps:
    r += 1
    wsr[f'B{r}'] = q
    wsr[f'B{r}'].font = F_BOLD
    wsr[f'B{r}'].alignment = WRAP
    wsr[f'C{r}'] = a
    wsr[f'C{r}'].font = F_BODY
    wsr[f'C{r}'].alignment = WRAP
    wsr.row_dimensions[r].height = 15 * max(2, math.ceil(len(a) / 90)) + 4
wsr.sheet_state = 'hidden'

# ----------------------------------------------------------------------------------
# 9. Sheet order, print setup, save
# ----------------------------------------------------------------------------------
order = ['Start here'] + [f'Level {i}' for i in range(1, 11)] + ['Words', ENG, 'Room notes', ANS]
wb._sheets = [wb[n] for n in order]
wb.active = 0
for ws in wb.worksheets:
    ws.sheet_view.zoomScale = 100
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.orientation = 'landscape'
    if ws.title != ENG:
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0
    ws.oddHeader.left.text = 'CRUK CONFIDENTIAL, INTERNAL USE ONLY'
    ws.oddFooter.left.text = 'The Econometrics Game v2'
    ws.oddFooter.right.text = 'Andrew Rajanathan, CRUK M&&D'
wb.save(OUT)
print('saved', OUT)
