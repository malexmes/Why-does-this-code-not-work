"""Build the Econometrics Game v3: a one-screen, app-style Excel game.

Usage: python build_app.py <wave1_results.pptx> <output_basename>
Writes <output_basename>.xlsx (formula-only twin). The VBA is added by add_vba.py.
"""
import sys, datetime, math
import numpy as np
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter, column_index_from_string
from openpyxl.chart import LineChart, BarChart, Reference
from openpyxl.chart.series import SeriesLabel
from openpyxl.chart.text import Text
from openpyxl.chart.data_source import StrRef
from openpyxl.chart.title import Title
from openpyxl.formatting.rule import FormulaRule, DataBarRule
from openpyxl.workbook.defined_name import DefinedName

PPTX = sys.argv[1]
OUT = sys.argv[2]

# ----------------------------------------------------------------------------------
# 1. AP's weekly series from the deck
# ----------------------------------------------------------------------------------
prs = Presentation(PPTX)
series = {}
dates = None
for idx in [16, 21, 23, 35, 46]:
    s = prs.slides[idx - 1]
    for sh in s.shapes:
        if not sh.has_chart:
            continue
        for plot in sh.chart.plots:
            cats = list(plot.categories)
            for ser in plot.series:
                series[(idx, ser.name.strip())] = np.array([0.0 if v is None else float(v) for v in ser.values])
            if idx == 23:
                dates = [datetime.date(1899, 12, 30) + datetime.timedelta(days=int(float(c))) for c in cats]
N = len(dates)
yrs = np.array([d.year for d in dates])
mons = np.array([d.month for d in dates])
S = lambda i, n: series[(i, n)]
actual = S(16, 'RFL Sign Ups')
cog_signups = S(16, 'COG Sign Ups')
consideration = S(46, 'Brand Consideration')
total_media = S(46, 'Total Media Spend')
dec = {n: v for (i, n), v in series.items() if i == 23}
sp = {n: v for (i, n), v in series.items() if i == 21}

ap_base = dec['Brand (Consideration)'] + dec['Emails']
ap_brake = dec['Price & Distribution'] + dec['Seasonality']
ap_sale = dec['Discount']
rfl_names = ['BRTV', 'DRTV', 'VOD', 'OOH', 'Radio', 'Regional', 'Paid Social', 'Paid Search', 'Digital Audio', 'Display',
             'YouTube', 'Press Partnership', 'Other Media']
ap_rfl = sum(dec[k] for k in rfl_names)
ap_halo = dec['Halo (COG)'] + dec['Halo (Legacy)'] + dec['Brand (TV, VOD, OOH)']
groups = [
    ('TV', ['BRTV RFL', 'DRTV RFL'], ['BRTV', 'DRTV']),
    ('Online video', ['VOD'], ['VOD']),
    ('Radio', ['Radio'], ['Radio']),
    ('Posters', ['OOH'], ['OOH']),
    ('Digital and social', ['Paid Social EMC', 'Paid Social Oliver', 'Google Ads Generic', 'Microsoft Ads Generic', 'Display',
                            'Digital Audio', 'YouTube', 'Regional', 'Press Partnership', 'Telemarketing Ethicall',
                            'Telemarketing SS', 'Direct Mail', 'Door Drops'],
     ['Paid Social', 'Paid Search', 'Display', 'Digital Audio', 'YouTube', 'Regional', 'Press Partnership', 'Other Media']),
]
ap_ch = [sum(dec[c] for c in ck) for _, _, ck in groups]
ch_spend = [sum(sp[s] for s in sk) / 1000 for _, sk, _ in groups]
rfl_spend = sum(sp[k] for k in sp if k not in ('BRTV Brand', 'VOD Brand')) / 1000
cog_spend = S(16, 'COG') / 1000
halo_spend = (S(16, 'COG') + S(16, 'Legacy') + S(16, 'Brand')) / 1000
ap_total = sum(dec.values())

ans_base = ap_base.mean()
ans_brake = -ap_brake.mean()
ans_sale = ap_sale[ap_sale > 0].mean()
shape_base = ap_base / ans_base
shape_brake = -ap_brake / ans_brake
shape_sale = ap_sale / ans_sale


def adstock(x, m):
    out = np.zeros_like(x); a = 0.0
    for i, v in enumerate(x):
        a = v + m * a; out[i] = a
    return out


def fit_mem(contrib, spend_k):
    best = None
    for m10 in range(0, 10):
        m = m10 / 10
        A = (1 - m) * adstock(spend_k, m)
        k = (A * contrib).sum() / (A * A).sum()
        sse = ((contrib - k * A) ** 2).sum()
        if best is None or sse < best[1]:
            best = (m10 * 10, sse)
    return best[0]


mem_rfl = fit_mem(ap_rfl, rfl_spend)
ans_rfl = float(ap_rfl.sum() / rfl_spend.sum())
ans_halo = float(ap_halo.sum() / halo_spend.sum())
ch_mem = [fit_mem(c, s) for c, s in zip(ap_ch, ch_spend)]
ch_ans = [float(c.sum() / s.sum()) for c, s in zip(ap_ch, ch_spend)]
print('answers: base %.0f brake %.0f sale %.0f rfl %.1f mem %d halo %.2f' % (ans_base, ans_brake, ans_sale, ans_rfl, mem_rfl, ans_halo))
print('channels:', [(g[0], round(k, 1), m) for g, k, m in zip(groups, ch_ans, ch_mem)])


def r2(model):
    return 1 - ((actual - model) ** 2).sum() / ((actual - actual.mean()) ** 2).sum()


b1 = ans_base * shape_base
b2 = -ans_brake * shape_brake
b3 = ans_sale * shape_sale
b4 = ans_rfl * (1 - mem_rfl / 100) * adstock(rfl_spend, mem_rfl / 100)
b5 = ans_halo * halo_spend
b6 = sum(k * (1 - m / 100) * adstock(s, m / 100) for k, m, s in zip(ch_ans, ch_mem, ch_spend))
m3 = b1 + b2 + b3; m4 = m3 + b4; m5 = m4 + b5; m6 = m3 + b6 + b5
print('answer R2: after sale %.2f, adverts %.2f, halo %.2f, split %.2f' % (r2(m3), r2(m4), r2(m5), r2(m6)))
jo24 = (yrs == 2024) & (mons <= 10)
jo25 = (yrs == 2025)
print('42%% cut at answer: %.0f (AP 9,139)' % (ans_halo * (cog_spend[jo24].sum() - cog_spend[jo25].sum())))

# ----------------------------------------------------------------------------------
# 2. Styles
# ----------------------------------------------------------------------------------
NAVY = '1B2A4A'; MAG = 'E60078'; CYAN = '009CEE'; GREY = '878787'; PALE = 'F4F6FA'
GREEN = '2E7D32'; GREENFILL = 'E6F4EA'; AMBER = 'C75000'; AMBERFILL = 'FFF4E5'; LIGHTBLUE = 'DDEBF7'
YEL = 'FFF2CC'
def F(size=11, bold=False, color='1A1A1A', italic=False):
    return Font(name='Arial', size=size, bold=bold, color=color, italic=italic)
FILL = lambda c: PatternFill('solid', fgColor=c)
thin = Side(style='thin', color='BBC7D9')
BOX = Border(left=thin, right=thin, top=thin, bottom=thin)
med = Side(style='medium', color=NAVY)
BTN_BORDER = Border(left=med, right=med, top=med, bottom=med)
CENTER = Alignment(horizontal='center', vertical='center', wrap_text=True)
LEFT = Alignment(horizontal='left', vertical='center', wrap_text=True)
TOP = Alignment(horizontal='left', vertical='top', wrap_text=True)

wb = Workbook()
wb.remove(wb.active)
ENG = 'Under the bonnet'; E = "'" + ENG + "'!"
STATE = 'State'; ST = STATE + '!'
STEPS = 'Steps'; SP = STEPS + '!'
ANS = 'Answers'; A_ = ANS + '!'
R0 = 6; R1 = R0 + N - 1
rng = lambda col: f"{E}${col}${R0}:${col}${R1}"
POUND = ST + '$B$48'

# State cell map
STEP = ST + '$B$1'
DIAL = {1: ST + '$B$2', 2: ST + '$B$3', 3: ST + '$B$4', 4: ST + '$B$5', 5: ST + '$B$6', 6: ST + '$B$7',
        8: ST + '$B$8', 9: ST + '$B$9', 10: ST + '$B$10', 11: ST + '$B$11', 12: ST + '$B$12'}
REVEAL = ST + '$B$14'; HINT = ST + '$B$15'
QCHOSEN = lambda q: f"{ST}$B${16 + q}"     # q = 1..10
STARS = ST + '$B$28'; STAGE = ST + '$B$30'; PROGRESS = ST + '$B$31'
H = {k: f"{ST}$B${v}" for k, v in dict(type=33, dialid=34, ans=35, tol=36, value=37, quiz=38, chosen=39, correct=40,
                                        verdict=41, r2=42, modelcol=43, bars=44, fmt=45, split_ok=46, tests_ok=47).items()}
BUD = lambda i: f"{ST}$B${50 + i}"        # 13 budget spends, i = 1..13
BUDMODE = ST + '$B$64'

# ----------------------------------------------------------------------------------
# 3. Engine
# ----------------------------------------------------------------------------------
eng = wb.create_sheet(ENG)
eng.sheet_properties.tabColor = GREY
eng['A1'] = 'Under the bonnet: every weekly number and formula the game uses. Nothing is hidden.'
eng['A1'].font = F(16, True, NAVY)
eng['A2'] = ("Columns A to D and the AP columns are Analytic Partners' Wave 1 weekly figures, read from the charts in their July 2026 "
             "results deck. Everything else is a formula pointing at the game's dials. Spend is in thousands of pounds.")
eng['A2'].font = F(11); eng.merge_cells('A2:T2'); eng.row_dimensions[2].height = 34
eng.freeze_panes = 'E6'
hdrs = [
    ('A', 'Week starting'), ('B', 'Year'), ('C', 'Month'), ('D', 'Real sign-ups'),
    ('E', 'Shape: base'), ('F', 'Shape: brake'), ('G', 'Shape: sale'),
    ('H', 'RFL adverts £k'), ('I', 'RFL adstock'), ('J', 'Committed Giving £k'), ('K', 'Other causes adverts £k'),
    ('N', 'TV £k'), ('O', 'Online video £k'), ('P', 'Radio £k'), ('Q', 'Posters £k'), ('R', 'Digital and social £k'),
    ('S', 'TV adstock'), ('T', 'Online video adstock'), ('U', 'Radio adstock'), ('V', 'Posters adstock'), ('W', 'Digital adstock'),
    ('X', 'Your block: base'), ('Y', 'Your block: brake'), ('Z', 'Your block: sale'), ('AA', 'Your block: RFL adverts'),
    ('AB', 'Your block: halo'), ('AE', 'Your block: TV'), ('AF', 'Your block: Online video'), ('AG', 'Your block: Radio'),
    ('AH', 'Your block: Posters'), ('AI', 'Your block: Digital and social'),
    ('AJ', 'Model 1 base'), ('AK', 'Model 2 brake'), ('AL', 'Model 3 sale'), ('AM', 'Model 4 adverts'), ('AN', 'Model 5 halo'),
    ('AO', 'Model 6 split'), ('AP', 'Your final model'),
    ('AQ', 'AP block: base'), ('AR', 'AP block: brake'), ('AS', 'AP block: sale'), ('AT', 'AP block: RFL adverts'),
    ('AU', 'AP block: halo'), ('AX', 'AP block: TV'), ('AY', 'AP block: Online video'), ('AZ', 'AP block: Radio'),
    ('BA', 'AP block: Posters'), ('BB', 'AP block: Digital and social'), ('BC', 'Wave 1 model'),
    ('BD', 'Wave 1 model 1'), ('BE', 'Wave 1 model 2'), ('BF', 'Wave 1 model 3'), ('BG', 'Wave 1 model 4'), ('BH', 'Wave 1 model 5'),
    ('BI', 'Wave 1 model 6'),
    ('BJ', 'Chart: Wave 1 answer'), ('BK', 'Chart: your model'), ('BL', 'Chart: adverts spend £k'),
    ('BQ', 'Miss'), ('BR', 'Miss squared'), ('BS', 'Miss change squared'), ('BT', 'Weekly % miss (weeks over 2,000)'),
    ('BU', 'Real new Committed Givers'), ('BV', 'Brand consideration'), ('BW', 'All CRUK media £k'),
]
for col, h in hdrs:
    c = eng[f'{col}5']; c.value = h; c.font = F(11, True, 'FFFFFF'); c.fill = FILL(NAVY); c.alignment = CENTER; c.border = BOX
    eng.column_dimensions[col].width = 13
eng.row_dimensions[5].height = 44
eng.column_dimensions['A'].width = 12
CL = {h: col for col, h in hdrs}
C = {h: column_index_from_string(col) for col, h in hdrs}
for i in range(N):
    r = R0 + i; p = r - 1
    eng[f'A{r}'] = dates[i]; eng[f'A{r}'].number_format = 'dd mmm yy'
    eng[f'B{r}'] = f'=YEAR(A{r})'; eng[f'C{r}'] = f'=MONTH(A{r})'
    eng[f'D{r}'] = round(float(actual[i]), 1)
    eng[f'E{r}'] = round(float(shape_base[i]), 4); eng[f'F{r}'] = round(float(shape_brake[i]), 4); eng[f'G{r}'] = round(float(shape_sale[i]), 4)
    eng[f'H{r}'] = round(float(rfl_spend[i]), 2)
    eng[f'I{r}'] = f'=H{r}' if i == 0 else f'=H{r}+{DIAL[5]}/100*I{p}'
    eng[f'J{r}'] = round(float(cog_spend[i]), 2)
    eng[f'K{r}'] = round(float(halo_spend[i]), 2)
    for j, colx in enumerate('NOPQR'):
        eng[f'{colx}{r}'] = round(float(ch_spend[j][i]), 2)
    for j, (colx, spcol) in enumerate(zip('STUVW', 'NOPQR')):
        eng[f'{colx}{r}'] = f'={spcol}{r}' if i == 0 else f'={spcol}{r}+{A_}$C${16 + j}/100*{colx}{p}'
    eng[f'X{r}'] = f'={DIAL[1]}*E{r}'
    eng[f'Y{r}'] = f'=-{DIAL[2]}*F{r}'
    eng[f'Z{r}'] = f'={DIAL[3]}*G{r}'
    eng[f'AA{r}'] = f'={DIAL[4]}*(1-{DIAL[5]}/100)*I{r}'
    eng[f'AB{r}'] = f'={DIAL[6]}*K{r}'
    for j, (colx, adcol) in enumerate(zip(['AE', 'AF', 'AG', 'AH', 'AI'], 'STUVW')):
        eng[f'{colx}{r}'] = f'={DIAL[8 + j]}*(1-{A_}$C${16 + j}/100)*{adcol}{r}'
    eng[f'AJ{r}'] = f'=X{r}'
    eng[f'AK{r}'] = f'=X{r}+Y{r}'
    eng[f'AL{r}'] = f'=AK{r}+Z{r}'
    eng[f'AM{r}'] = f'=AL{r}+AA{r}'
    eng[f'AN{r}'] = f'=AM{r}+AB{r}'
    eng[f'AO{r}'] = f'=AL{r}+AE{r}+AF{r}+AG{r}+AH{r}+AI{r}+AB{r}'
    eng[f'AP{r}'] = f'=IF({H["split_ok"]}=1,AO{r},AN{r})'
    eng[f'AQ{r}'] = round(float(ap_base[i]), 1); eng[f'AR{r}'] = round(float(ap_brake[i]), 1); eng[f'AS{r}'] = round(float(ap_sale[i]), 1)
    eng[f'AT{r}'] = round(float(ap_rfl[i]), 1); eng[f'AU{r}'] = round(float(ap_halo[i]), 1)
    for j, colx in enumerate(['AX', 'AY', 'AZ', 'BA', 'BB']):
        eng[f'{colx}{r}'] = round(float(ap_ch[j][i]), 1)
    eng[f'BC{r}'] = f'=AQ{r}+AR{r}+AS{r}+AT{r}+AU{r}'
    eng[f'BD{r}'] = f'=AQ{r}'; eng[f'BE{r}'] = f'=AQ{r}+AR{r}'; eng[f'BF{r}'] = f'=BE{r}+AS{r}'
    eng[f'BG{r}'] = f'=BF{r}+AT{r}'; eng[f'BH{r}'] = f'=BG{r}+AU{r}'; eng[f'BI{r}'] = f'=BC{r}'
    eng[f'BK{r}'] = f'=CHOOSE({H["modelcol"]}+1,D{r},AJ{r},AK{r},AL{r},AM{r},AN{r},AO{r})'
    eng[f'BJ{r}'] = f'=IF({REVEAL}=1,CHOOSE({H["modelcol"]}+1,D{r},BD{r},BE{r},BF{r},BG{r},BH{r},BI{r}),BK{r})'
    eng[f'BL{r}'] = f'=IF({H["bars"]}=1,H{r},0)'
    eng[f'BQ{r}'] = f'=D{r}-AP{r}'; eng[f'BR{r}'] = f'=BQ{r}^2'
    eng[f'BS{r}'] = 0 if i == 0 else f'=(BQ{r}-BQ{p})^2'
    eng[f'BT{r}'] = f'=IF(D{r}>2000,ABS(BQ{r})/D{r},"")'
    eng[f'BU{r}'] = round(float(cog_signups[i]), 0)
    eng[f'BV{r}'] = round(float(consideration[i]), 4); eng[f'BV{r}'].number_format = '0.0%'
    eng[f'BW{r}'] = round(float(total_media[i]) / 1000, 1)
    for col, _ in hdrs:
        cc = eng[f'{col}{r}']; cc.font = F(9)
        if col not in ('A', 'B', 'C', 'BV'):
            cc.number_format = '#,##0.0' if col in ('E', 'F', 'G') else '#,##0'
ACT = rng('D'); FINAL = rng('AP'); YEAR = rng('B'); MONTH = rng('C'); YOURS = rng('BK')
JO24 = lambda col: f'SUMIFS({col},{YEAR},2024,{MONTH},"<=10")'
JO25 = lambda col: f'SUMIFS({col},{YEAR},2025,{MONTH},"<=10")'

# ----------------------------------------------------------------------------------
# 4. Answers (very hidden)
# ----------------------------------------------------------------------------------
ans = wb.create_sheet(ANS)
ans['A1'] = 'Answer key, fitted to the Wave 1 decomposition. Hidden so the game stays a game.'
for r, lab, v in [(3, 'Base: sign-ups a week', round(ans_base)), (4, 'Brake: sign-ups lost a week', round(ans_brake)),
                  (5, 'Sale: extra sign-ups in a sale week', round(ans_sale, -1)), (6, 'RFL adverts: sign-ups per £1,000', round(ans_rfl, 1)),
                  (7, 'RFL adverts: memory %', mem_rfl), (8, 'Halo: sign-ups per £1,000 of other adverts', round(ans_halo, 2)),
                  (11, 'Tolerance on a dial', 0.25), (12, 'Tolerance on memory (points)', 10), (13, 'Tolerance on a channel', 0.35)]:
    ans[f'B{r}'] = lab; ans[f'C{r}'] = v
for j, ((g, _, _), m, k) in enumerate(zip(groups, ch_mem, ch_ans)):
    ans[f'B{16 + j}'] = f'{g}: memory %'; ans[f'C{16 + j}'] = m
    ans[f'D{16 + j}'] = f'{g}: sign-ups per £1,000'; ans[f'E{16 + j}'] = round(k, 1)
ans.column_dimensions['B'].width = 44; ans.column_dimensions['D'].width = 40
ans.sheet_state = 'veryHidden'

# ----------------------------------------------------------------------------------
# 5. Steps content
# ----------------------------------------------------------------------------------
steps = wb.create_sheet(STEPS)
cols = ['step', 'chapter', 'chapter name', 'title', 'sentence', 'type', 'dial id', 'small step', 'big step', 'answer', 'tolerance',
        'tick', 'hint', 'badge', 'option A', 'option B', 'option C', 'correct', 'quiz index', 'model col', 'show bars', 'unit', 'format']
for j, h in enumerate(cols):
    steps.cell(row=1, column=1 + j, value=h).font = F(10, True)
R2_OF = lambda col: f'MAX(0,1-SUMXMY2({ACT},{rng(col)})/DEVSQ({ACT}))'
LOST = f'{DIAL[6]}*({JO24(rng("J"))}-{JO25(rng("J"))})'
content = [
    # step, ch, chname, title, sentence, type, dial, small, big, answer, tol, tick, hint, badge, A, B, C, correct, quizidx, modelcol, bars, unit, fmt
    (1, 1, 'The receipts', 'Three years of receipts',
     'Every week for nearly three years, this many people signed up for Race for Life. Look at the shape.',
     'info', 0, 0, 0, 0, 0, 'The tall spikes come in spring. From August to December it goes quiet. Something causes every spike.', '', '', '', '', '', 0, 0, 0, 0, '', ''),
    (2, 1, 'The receipts', 'What made the tall spikes?', 'Tap the answer you think is right.',
     'quiz', 0, 0, 0, 0, 0, 'Yes. Every spike is a discount week. AP tested weather and bank holidays and found nothing.',
     'Look at the spikes. They come in the same weeks each year, whatever the weather.', 'Driver',
     'Sale weeks', 'Warm weather', 'Bank holidays', 1, 1, 0, 0, '', ''),
    (3, 1, 'The receipts', 'You are the detective',
     'A detective with three years of receipts works out what caused each sign-up. That is all econometrics is. You are about to do it, one block at a time.',
     'info', 0, 0, 0, 0, 0, 'Analytic Partners did this for us in 2026 and called it Wave 1. Every number in this game is theirs.', '', 'Econometrics', '', '', '', 0, 0, 0, 0, '', ''),
    (4, 2, 'The shop', 'Start with the quiet weeks',
     'In autumn almost nobody signs up. Almost. Tap + until your blue line sits on the quiet weeks.',
     'dial', 1, 100, 500, f'={A_}C3', f'={A_}C11', 'About 900 a week sign up whatever we do. People who already know us, and the ones we email.',
     'Look at August to December. Hundreds, not thousands.', 'Baseline', '', '', '', 0, 0, 1, 0, 'sign-ups a week', '#,##0'),
    (5, 2, 'The shop', 'If every advert and every sale stopped, how many would sign up a week?', 'Tap one.',
     'quiz', 0, 0, 0, 0, 0, 'Right. The base is small. Almost every other sign-up was caused by something we did.',
     'Look at where your blue line sat in the last step.', 'Baseline', 'About 900', 'About 9,000', 'Nobody at all', 1, 2, 1, 0, '', ''),
    (6, 2, 'The shop', 'Now the brake',
     'A dearer ticket and a shut window push people away. Tap + until the brake feels right. Your line dips below zero in autumn. That is fine for now.',
     'dial', 2, 100, 500, f'={A_}C4', f'={A_}C11', 'About 950 a week are lost to price and the shut window. Blocks can pull down as well as push up.',
     'About the same size as the base, pulling the other way.', 'Negative driver', '', '', '', 0, 0, 2, 0, 'sign-ups lost a week', '#,##0'),
    (7, 3, 'The sale', 'Every spike is a sale', 'Tap + until the spikes on your blue line reach the pink ones.',
     'dial', 3, 500, 2000, f'={A_}C5', f'={A_}C11', 'A sale week adds about 7,600 sign-ups. Sales explained a fifth of 2025 sign-ups.',
     'How tall is a spike above the weeks around it? Thousands.', 'Discount driver', '', '', '', 0, 0, 3, 0, 'extra sign-ups in a sale week', '#,##0'),
    (8, 3, 'The sale', 'Look at the week after a sale. What happens?', 'Find a spike on the pink line and look one week to the right.',
     'quiz', 0, 0, 0, 0, 0, 'Yes. Some people signed up early to get the deal. The model counts the dip so the sale is not over-credited.',
     'The week after a sale is lower than you would expect.', 'Pull-forward', 'A dip: some people signed up early', 'An even bigger spike', 'Nothing changes', 1, 3, 3, 0, '', ''),
    (9, 3, 'The sale', 'Half the story',
     f'="Your model now explains "&TEXT({R2_OF("AL")},"0%")&" of the weekly pattern. Sales and the shop are half the story. Adverts are the other half."',
     'info', 0, 0, 0, 0, 0, 'A model with no adverts in it cannot say what adverts did. That is the next chapter.', '', 'Pattern explained', '', '', '', 0, 0, 3, 0, '', ''),
    (10, 4, 'The kettle', 'Now the adverts',
     'The grey bars are what we spent on Race for Life adverts each week. Tap + to set how many sign-ups each £1,000 brings.',
     'dial', 4, 1, 5, f'={A_}C6', f'={A_}C11',
     f'="About 16 sign-ups per £1,000. That is £"&TEXT({DIAL[4]}*{POUND}/1000,"0.00")&" of sponsorship for every £1 spent."',
     'Somewhere between 5 and 30.', 'Response', '', '', '', 0, 0, 4, 1, 'sign-ups per £1,000', '#,##0.0'),
    (11, 4, 'The kettle', 'The kettle',
     'Adverts keep working after they stop, like a kettle stays hot. Tap + on the memory until your line matches the shape of the pink one.',
     'dial', 5, 10, 10, f'={A_}C7', f'={A_}C12', 'Memory of 80%: the effect fades slowly over months. Judge a campaign in week one and you miss most of it.',
     'Try a high number. Adverts fade over months, not days.', 'Adstock', '', '', '', 0, 0, 4, 1, "% of last week's effect still here", '0"%"'),
    (12, 4, 'The kettle', 'Every £1 of Race for Life adverts brought back how much sponsorship in 2025?', 'Tap one.',
     'quiz', 0, 0, 0, 0, 0, 'Yes. £1.66 on average. Remember the word average. It matters in chapter 8.',
     'Look at the £ figure from two steps ago.', 'ROI, return per £1', '£1.66', '16p', '£16', 1, 4, 4, 1, '', ''),
    (13, 5, 'The next field', 'Rain on the next field',
     'Committed Giving, Legacy and Brand adverts never mention Race for Life. Yet sign-ups rose when they ran. Tap + to add their effect.',
     'dial', 6, 0.5, 2, f'={A_}C8', f'={A_}C11', 'About 3 per £1,000. Small per pound, but the spend is large. It adds up to 15% of sign-ups.',
     'Small. Under 10.', 'Halo', '', '', '', 0, 0, 5, 0, 'sign-ups per £1,000 of other adverts', '0.0'),
    (14, 5, 'The next field', 'What the 2025 cut cost',
     f'="In 2025 we spent 42% less on Committed Giving adverts. Your model says Race for Life lost "&TEXT({LOST},"#,##0")&" sign-ups because of it. Wave 1 said 9,139."',
     'info', 0, 0, 0, 0, 0, 'Two budgets, one shared result. Yours is lower than Wave 1\'s because you gave every other cause the same strength per pound. AP found Committed Giving the strongest of the three.', '', 'Halo', '', '', '', 0, 0, 5, 0, '', ''),
    (15, 5, 'The next field', 'Adverts for one cause helping another is called?', 'Tap one.',
     'quiz', 0, 0, 0, 0, 0, 'A halo. Committed Giving adverts were quietly filling Race for Life too.',
     'Think of rain on the next field.', 'Halo', 'A halo', 'A discount', 'Seasonality', 1, 5, 5, 0, '', ''),
    (16, 6, 'Who sang louder', 'Split the adverts into five',
     'TV, video, radio, posters, digital. They mostly run in the same weeks. Tap − and + on each until all five ticks show.',
     'dials5', 0, 1, 1, 0, f'={A_}C13', 'All five about right. Notice how loose the checks were. That is the point.',
     'Your chapter 4 number is the guide for each. Posters are much weaker.', 'Multicollinearity', '', '', '', 0, 0, 6, 0, '', ''),
    (17, 6, 'Who sang louder', 'The group is firm, the split is soft',
     f'="With one adverts block your pattern score was "&TEXT({R2_OF("AN")},"0%")&". With five it is "&TEXT({R2_OF("AO")},"0%")&". The split changes the story, not the fit."',
     'info', 0, 0, 0, 0, 0, 'Two people singing the same note: you hear the sound but cannot tell who is louder. Quote the range on any single channel.', '', 'Multicollinearity', '', '', '', 0, 0, 6, 0, '', ''),
    (18, 6, 'Who sang louder', 'Which channel was AP least sure about?', 'Tap one.',
     'quiz', 0, 0, 0, 0, 0, 'Posters: about 60p back per £1, with a wide range either side.',
     'Which one needed the smallest number?', 'Confidence interval', 'Posters', 'TV', 'Digital and social', 1, 6, 6, 0, '', ''),
    (19, 7, 'Any good?', 'Is your model any good?', 'Four questions judge any model. Your scores sit beside Wave 1\'s.',
     'tests', 0, 0, 0, 0, 0, 'Green all round. Your model explains the pattern, predicts weeks it never saw and leaves little unexplained.',
     'Amber? Go back a chapter and adjust that block.', 'R squared, MAPE, holdout', '', '', '', 0, 0, 6, 0, '', ''),
    (20, 7, 'Any good?', 'A model copies the past perfectly but predicts 2025 badly. What went wrong?', 'Tap one.',
     'quiz', 0, 0, 0, 0, 0, 'Yes. That is why AP test every model on weeks it never saw.',
     'Think of a student who memorises last year\'s exam paper.', 'Holdout test', 'It learned the noise, not the pattern', 'It needs more adverts', '2025 was just unusual', 1, 7, 6, 0, '', ''),
    (21, 7, 'Any good?', 'You built the model', 'Seven chapters, one model, every number Wave 1\'s. Next: where should the next pound go?',
     'info', 0, 0, 0, 0, 0, 'Say it in the room: about half of Race for Life sign-ups come from Race for Life adverts, a fifth from sales and 15% from other causes\' adverts.', '', 'Decomposition', '', '', '', 0, 0, 6, 0, '', ''),
]
for row in content:
    for j, v in enumerate(row):
        steps.cell(row=row[0] + 1, column=1 + j, value=v)
steps.sheet_state = 'hidden'
NSTEP = len(content)
IDX = lambda col: f'INDEX({SP}${col}$2:${col}${NSTEP + 1},{STEP})'

# ----------------------------------------------------------------------------------
# 6. State sheet
# ----------------------------------------------------------------------------------
st = wb.create_sheet(STATE)
st.column_dimensions['A'].width = 44; st.column_dimensions['B'].width = 14
st['A1'] = 'Step (1 to 21)'; st['B1'] = 1
labels = {2: 'Base: sign-ups a week', 3: 'Brake: sign-ups lost a week', 4: 'Sale: extra sign-ups', 5: 'Adverts: sign-ups per £1,000',
          6: 'Memory %', 7: 'Halo: sign-ups per £1,000', 8: 'TV per £1,000', 9: 'Online video per £1,000', 10: 'Radio per £1,000',
          11: 'Posters per £1,000', 12: 'Digital and social per £1,000', 14: 'Show answer (1 = yes)', 15: 'Show hint (1 = yes)'}
for r, lab in labels.items():
    st[f'A{r}'] = lab; st[f'B{r}'] = 0
for q in range(1, 11):
    st[f'A{16 + q}'] = f'Quiz {q} answer chosen (1 to 3)'; st[f'B{16 + q}'] = 0; st[f'C{16 + q}'] = 1
st['A28'] = 'Stars'; st['B28'] = '=SUMPRODUCT(--(B17:B26=C17:C26))'
st['A30'] = 'Stage (1 Play, 2 Budget, 3 Long game, 4 Giving, 5 Finish)'; st['B30'] = 1
st['A31'] = 'Progress (of 24)'; st['B31'] = '=IF(B30=1,B1,MIN(24,21+B30-1))'
helpers = {
    33: ('Step type', f'={IDX("F")}'), 34: ('Dial id', f'={IDX("G")}'), 35: ('Answer', f'={IDX("J")}'), 36: ('Tolerance', f'={IDX("K")}'),
    37: ('Current value', f'=IF(B34=0,0,INDEX($B$2:$B$12,B34))'), 38: ('Quiz index', f'={IDX("S")}'),
    39: ('Chosen', '=IF(B38=0,0,INDEX($B$17:$B$26,B38))'), 40: ('Correct', f'={IDX("R")}'),
    41: ('Dial verdict (0 low, 1 ok, 2 high)', '=IF(B34=0,1,IF(B34=5,IF(B37<B35-B36,0,IF(B37>B35+B36,2,1)),IF(B37<B35*(1-B36),0,IF(B37>B35*(1+B36),2,1))))'),
    42: ('Pattern explained (chart model)', f'=MAX(0,1-SUMXMY2({ACT},{YOURS})/DEVSQ({ACT}))'),
    43: ('Model column', f'={IDX("T")}'), 44: ('Show bars', f'={IDX("U")}'), 45: ('Number format', f'={IDX("W")}'),
    46: ('Split typed (all five above zero)', '=IF(MIN(B8:B12)>0,1,0)'),
    47: ('Tests all green', '=IF(AND(B70=1,B71=1,B72=1,B73=1),1,0)'),
    48: ('Sponsorship per sign-up £', 90),
}
for r, (lab, f) in helpers.items():
    st[f'A{r}'] = lab; st[f'B{r}'] = f
# budget
channels = [('DRTV', 984684, 1.955, 1375338), ('Online video', 899101, 1.861, 768640), ('Radio', 834255, 1.753, 666525),
            ('Paid social', 1062828, 1.591, 1077485), ('Regional', 748218, 1.710, 460481), ('Posters', 521408, 0.538, 320322),
            ('Digital audio', 458113, 1.708, 657756), ('Display', 441686, 1.740, 552287), ('Paid search', 317693, 1.7, 374699),
            ('YouTube', 120519, 1.850, 150025), ('Telemarketing', 102187, 1.2, 71033), ('Direct mail', 40000, 1.9, 57392),
            ('Door drops', 22184, 2.0, 20892)]
st['A50'] = 'Budget: your spend by channel, £k'
for i, (nm, spend, roi, apv) in enumerate(channels, start=1):
    st[f'A{50 + i}'] = nm; st[f'B{50 + i}'] = round(spend / 1000); st[f'C{50 + i}'] = round(spend / 1000)
    st[f'D{50 + i}'] = roi; st[f'E{50 + i}'] = round(apv / 1000)
st['A64'] = 'Budget plan mode (0 mine, 1 AP)'; st['B64'] = 0
st['C50'] = '2025 £k'; st['D50'] = 'Wave 1 return'; st['E50'] = "AP plan £k"
# tests
tests = [
    (70, 'Pattern explained', f'={R2_OF("AP")}', 0.94, '=IF(B75>=0.85,1,0)', '0%'),
    (71, 'Typical weekly miss', f'=IFERROR(AVERAGE({rng("BT")}),1)', 0.28, '=IF(B76<=0.35,1,0)', '0%'),
    (72, 'Predicts 2025', f'=MAX(0,1-SUMPRODUCT(({YEAR}=2025)*({ACT}-{FINAL})^2)/SUMPRODUCT(({YEAR}=2025)*({ACT}-SUMPRODUCT(({YEAR}=2025)*{ACT})/COUNTIF({YEAR},2025))^2))', None, '=IF(B77>=0.8,1,0)', '0%'),
    (73, 'Unexplained 2025 change', f'=IFERROR((({JO25(ACT)}-{JO24(ACT)})-({JO25(FINAL)}-{JO24(FINAL)}))/{JO24(ACT)},1)', 0.014, '=IF(ABS(B78)<=0.05,1,0)', '0.0%'),
]
st['A69'] = 'Tests: green flag'; st['A74'] = 'Tests: your value'
for r, lab, f, w1, flag, fmt in tests:
    st[f'A{r}'] = lab + ' ok'; st[f'B{r}'] = flag
    st[f'A{r + 5}'] = lab; st[f'B{r + 5}'] = f; st[f'B{r + 5}'].number_format = fmt
    st[f'C{r + 5}'] = w1 if w1 is not None else 'n/a'
    if w1 is not None:
        st[f'C{r + 5}'].number_format = fmt
for r in range(1, 80):
    st[f'A{r}'].font = F(10); st[f'B{r}'].font = F(10)
for r in list(range(1, 13)) + [14, 15] + list(range(17, 27)) + [48] + list(range(51, 64)) + [64]:
    st[f'B{r}'].fill = FILL(YEL); st[f'B{r}'].border = BOX
st['D1'] = ('No macros? This tab is the control panel. Type the step number in B1 (1 to 21), the dials in B2 to B12, '
            'quiz answers in B17 to B26 and the budget in B51 to B63. Set B30 to 2, 3, 4 or 5 to move to the later chapters.')
st['D1'].font = F(10, italic=True); st['D1'].alignment = TOP; st.merge_cells('D1:K6')
st.sheet_state = 'hidden'

# ----------------------------------------------------------------------------------
# 7. App sheets
# ----------------------------------------------------------------------------------
def app_frame(ws, tab_color, chapter_formula, title_formula, sentence_formula, back_label='◀  Back', next_label='Next  ▶'):
    ws.sheet_properties.tabColor = tab_color
    ws.sheet_view.showGridLines = False
    ws.sheet_view.zoomScale = 100
    widths = {'A': 2}
    for c in 'BCDEFGHIJKLMNOPQR':
        widths[c] = 5.2
    widths.update({'S': 2, 'T': 40, 'U': 5.5, 'V': 5.5, 'W': 5.5, 'X': 7, 'Y': 7, 'Z': 7, 'AA': 2})
    for c, w in widths.items():
        ws.column_dimensions[c].width = w
    ws.row_dimensions[1].height = 8
    ws.row_dimensions[2].height = 20
    ws.row_dimensions[3].height = 7
    ws.row_dimensions[4].height = 8
    ws.row_dimensions[5].height = 40
    ws.row_dimensions[6].height = 6
    ws.row_dimensions[7].height = 48
    ws.row_dimensions[8].height = 8
    ws.row_dimensions[9].height = 10
    # chapter label
    ws['B2'] = chapter_formula; ws['B2'].font = F(11, True, GREY); ws.merge_cells('B2:M2')
    ws['T2'] = f'=REPT("★",{STARS})&REPT("☆",10-{STARS})'; ws['T2'].font = F(14, False, MAG)
    ws['T2'].alignment = Alignment(horizontal='right', vertical='center'); ws.merge_cells('T2:Z2')
    # progress bar: 24 cells B3:Y3
    for k in range(24):
        c = ws.cell(row=3, column=2 + k); c.fill = FILL('E9EDF3')
    ws.conditional_formatting.add('B3:Y3', FormulaRule(formula=[f'COLUMN()-1<={PROGRESS}'], fill=FILL(MAG)))
    ws['B5'] = title_formula; ws['B5'].font = F(20, True, NAVY); ws['B5'].alignment = LEFT; ws.merge_cells('B5:Z5')
    ws['B7'] = sentence_formula; ws['B7'].font = F(13); ws['B7'].alignment = TOP; ws.merge_cells('B7:Z7')
    # buttons row 29
    ws.row_dimensions[28].height = 10
    ws.row_dimensions[29].height = 34
    def btn(rng_, text, fill, fc='FFFFFF', size=13):
        ws.merge_cells(rng_)
        c = ws[rng_.split(':')[0]]; c.value = text; c.font = F(size, True, fc); c.fill = FILL(fill); c.alignment = CENTER
        for row in ws[rng_]:
            for cell in row:
                cell.border = BOX
    btn('B29:E29', back_label, 'E9EDF3', NAVY)
    btn('T29:Z29', next_label, MAG)
    ws.freeze_panes = None
    return btn


# ---- Play
play = wb.create_sheet('Play')
play.sheet_properties.codeName = 'Sheet1'
btn = app_frame(play, MAG, f'="Chapter "&{IDX("B")}&"  ·  "&{IDX("C")}', f'={IDX("D")}', f'={IDX("E")}')
btn('G29:J29', f'=IF({HINT}=1,"Hide hint","Hint")', 'E9EDF3', NAVY)
btn('L29:O29', f'=IF({REVEAL}=1,"Hide answer","Show answer")', 'E9EDF3', NAVY)
btn('Q29:R29', '↺', 'E9EDF3', NAVY)
play['Q31'] = 'Start again'; play['Q31'].font = F(8, color=GREY); play.merge_cells('Q31:R31')
play['Q31'].alignment = Alignment(horizontal='center')
TYPE = H['type']; VAL = H['value']; ANSV = H['ans']; VERD = H['verdict']; CHO = H['chosen']; COR = H['correct']
for r in range(10, 28):
    play.row_dimensions[r].height = 26
play.row_dimensions[16].height = 34
play.row_dimensions[18].height = 60
play.row_dimensions[23].height = 24
play.row_dimensions[25].height = 24
play.row_dimensions[27].height = 24
# slot rows 11 to 15
names5 = [g[0] for g in groups]
test_names = ['Pattern explained', 'Typical weekly miss', 'Predicts 2025', 'Unexplained 2025 change']
for k in range(5):
    r = 11 + k
    quiz_opt = {11: 'A', 13: 'B', 15: 'C'}.get(r)
    parts = []
    if r == 11:
        parts.append(f'IF({TYPE}="dial",TEXT({VAL},{H["fmt"]}),')
    if quiz_opt:
        col = {'A': 'O', 'B': 'P', 'C': 'Q'}[quiz_opt]
        parts.append(f'IF({TYPE}="quiz","{quiz_opt}.   "&{IDX(col)},')
    parts.append(f'IF({TYPE}="dials5","{names5[k]}:  "&TEXT({DIAL[8 + k]},"0.0")&" per £1,000",')
    if k < 4:
        parts.append(f'IF({TYPE}="tests","{test_names[k]}",')
    parts.append('""')
    play[f'T{r}'] = '=' + ''.join(parts) + ')' * (len(parts) - 1)
    play[f'T{r}'].alignment = LEFT
    play[f'T{r}'].font = F(12)
    # dials5 buttons and verdicts
    play[f'U{r}'] = f'=IF({TYPE}="dials5","−","")'; play[f'V{r}'] = f'=IF({TYPE}="dials5","+","")'
    tol = A_ + '$C$13'; a = A_ + f'$E${16 + k}'; d = DIAL[8 + k]
    play[f'W{r}'] = f'=IF({TYPE}="dials5",IF({d}<{a}*(1-{tol}),"↑",IF({d}>{a}*(1+{tol}),"↓","✓")),"")'
    for cc in 'UV':
        play[f'{cc}{r}'].font = F(14, True, NAVY); play[f'{cc}{r}'].alignment = CENTER
    play[f'W{r}'].font = F(14, True, GREEN); play[f'W{r}'].alignment = CENTER
    if k < 4:
        yours = f'{ST}$B${75 + k}'; w1 = f'{ST}$C${75 + k}'; ok = f'{ST}$B${70 + k}'
        fmt = tests[k][5]
        play[f'X{r}'] = f'=IF({TYPE}="tests",TEXT({yours},"{fmt}"),"")'
        play[f'Y{r}'] = f'=IF({TYPE}="tests",IF(ISNUMBER({w1}),TEXT({w1},"{fmt}"),"–"),"")'
        play[f'Z{r}'] = f'=IF({TYPE}="tests",IF({ok}=1,"●","●"),"")'
        play[f'X{r}'].font = F(12, True); play[f'Y{r}'].font = F(12, color=GREY)
        play[f'X{r}'].alignment = CENTER; play[f'Y{r}'].alignment = CENTER; play[f'Z{r}'].alignment = CENTER
        play[f'Z{r}'].font = F(16, True, GREEN)
        play.conditional_formatting.add(f'Z{r}', FormulaRule(formula=[f'AND({TYPE}="tests",{ok}=0)'], font=Font(name='Arial', size=16, bold=True, color=AMBER)))
play['T11'].font = F(12)
play.conditional_formatting.add('T11', FormulaRule(formula=[f'{TYPE}="dial"'], font=Font(name='Arial', size=26, bold=True, color=NAVY)))
# tests header row 10
play['X10'] = f'=IF({TYPE}="tests","Yours","")'; play['Y10'] = f'=IF({TYPE}="tests","Wave 1","")'
play['X10'].font = F(9, color=GREY); play['Y10'].font = F(9, color=GREY); play['X10'].alignment = CENTER; play['Y10'].alignment = CENTER
play['T10'] = f'=IF({TYPE}="dial","Your number",IF({TYPE}="quiz","Tap one",IF({TYPE}="dials5","Sign-ups per £1,000",IF({TYPE}="tests","Four tests",""))))'
play['T10'].font = F(9, color=GREY)
# quiz option styling and chosen highlight
for r, k in ((11, 1), (13, 2), (15, 3)):
    play.conditional_formatting.add(f'T{r}', FormulaRule(formula=[f'{TYPE}="quiz"'], border=BOX, fill=FILL(PALE)))
    play.conditional_formatting.add(f'T{r}', FormulaRule(formula=[f'AND({TYPE}="quiz",{CHO}={k})'], fill=FILL(LIGHTBLUE), font=Font(name='Arial', size=12, bold=True, color=NAVY)))
# dial buttons row 16
play['T16'] = f'=IF({TYPE}="dial",{IDX("V")},"")'; play['T16'].font = F(10, color=GREY); play['T16'].alignment = LEFT
for cc, txt in (('U', '−−'), ('V', '−'), ('X', '+'), ('Y', '++')):
    play[f'{cc}16'] = f'=IF({TYPE}="dial","{txt}","")'
    play[f'{cc}16'].font = F(14, True, 'FFFFFF'); play[f'{cc}16'].alignment = CENTER
    play.conditional_formatting.add(f'{cc}16', FormulaRule(formula=[f'{TYPE}="dial"'], fill=FILL(NAVY)))
# feedback T18:Z21
play.merge_cells('T18:Z21')
tick = IDX('L'); hint = IDX('M'); nudge = IDX('M')
fb = (f'=IF({TYPE}="dial",IF({VAL}=0,"Tap + to start."&IF({HINT}=1," "&{hint},""),'
      f'IF({VERD}=0,"↑ Higher."&IF({HINT}=1," "&{hint},""),IF({VERD}=2,"↓ Lower."&IF({HINT}=1," "&{hint},""),"✓ "&{tick}))),'
      f'IF({TYPE}="quiz",IF({CHO}=0,"Pick one.",IF({CHO}={COR},"✓ "&{tick},"✗ Not quite. "&{nudge})),'
      f'IF({TYPE}="dials5",IF(AND(W11="✓",W12="✓",W13="✓",W14="✓",W15="✓"),"✓ "&{tick},IF({HINT}=1,{hint},"Tap − and + on each channel until five ticks show.")),'
      f'IF({TYPE}="tests",IF({H["tests_ok"]}=1,"✓ "&{tick},{hint}),{tick}))))')
play['T18'] = fb; play['T18'].font = F(12); play['T18'].alignment = TOP
play.conditional_formatting.add('T18', FormulaRule(formula=['LEFT(T18,1)="✓"'], fill=FILL(GREENFILL), font=Font(name='Arial', size=12, color=GREEN)))
play.conditional_formatting.add('T18', FormulaRule(formula=['OR(LEFT(T18,1)="↑",LEFT(T18,1)="↓",LEFT(T18,1)="✗")'], fill=FILL(AMBERFILL), font=Font(name='Arial', size=12, color=AMBER)))
play.conditional_formatting.add('T18', FormulaRule(formula=[f'{TYPE}="info"'], fill=FILL(PALE)))
play.merge_cells('T23:Z23')
play['T23'] = f'=IF(OR(LEFT(T18,1)="✓",{TYPE}="info"),IF({IDX("N")}="","","Economists call this: "&{IDX("N")}),"")'
play['T23'].font = F(10, True, STEELBLUE := '4A6FA5'); play['T23'].alignment = LEFT
play.merge_cells('T25:Z25')
play['T25'] = f'=IF({H["modelcol"]}>=3,"Pattern explained so far: "&TEXT({H["r2"]},"0%"),"")'
play['T25'].font = F(10, color=GREY); play['T25'].alignment = LEFT
play.merge_cells('T27:Z27')
play['T27'] = (f'=IF({REVEAL}=1,IF({TYPE}="dial","Wave 1\'s number: "&TEXT({ANSV},{H["fmt"]}),'
               f'IF({TYPE}="dials5","Wave 1: "&TEXT({A_}E16,"0.0")&", "&TEXT({A_}E17,"0.0")&", "&TEXT({A_}E18,"0.0")&", "&TEXT({A_}E19,"0.0")&", "&TEXT({A_}E20,"0.0"),'
               f'IF({TYPE}="quiz","Wave 1\'s answer: "&CHOOSE({COR},"A","B","C"),"Wave 1\'s model is the grey dashed line."))),"")')
play['T27'].font = F(10, True, MAG); play['T27'].alignment = LEFT
play['B31'] = f'=IF({TYPE}="dial","Tip: −− and ++ move in bigger steps.","")'
play['B31'].font = F(9, color=GREY); play.merge_cells('B31:M31')
# chart
play['B9'] = f'="Real sign-ups (pink) and your model (blue), "&IF({H["bars"]}=1,"with adverts spend as grey bars","week by week")'
play['B9'].font = F(9, color=GREY)
bar = BarChart(); bar.type = 'col'; bar.style = 2; bar.width = 19.5; bar.height = 12.4
bar.legend.position = 'b'; bar.y_axis.majorGridlines = None
bar.y_axis.number_format = '#,##0'; bar.y_axis.title = None
bar.x_axis.number_format = 'mmm yy'; bar.x_axis.delete = False; bar.y_axis.delete = False
bar.x_axis.tickLblSkip = 13; bar.gapWidth = 30
bar.add_data(Reference(eng, min_col=C['Chart: adverts spend £k'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
bar.series[0].graphicalProperties.solidFill = 'D9D9D9'; bar.series[0].graphicalProperties.line.solidFill = 'D9D9D9'
bar.series[0].tx = SeriesLabel(v='Adverts spend £k')
bar.y_axis.axId = 200; bar.y_axis.crosses = 'max'; bar.y_axis.delete = True; bar.y_axis.scaling.min = 0; bar.y_axis.scaling.max = 1200
bar.set_categories(Reference(eng, min_col=1, min_row=R0, max_row=R1))
line = LineChart()
for colname, label, color, dash in (('Chart: Wave 1 answer', 'Wave 1 answer', '9E9E9E', 'dash'),
                                    ('Chart: your model', 'Your model', CYAN, None), ('Real sign-ups', 'Real sign-ups', MAG, None)):
    line.add_data(Reference(eng, min_col=C[colname], min_row=R0 - 1, max_row=R1), titles_from_data=True)
    s = line.series[-1]; s.tx = SeriesLabel(v=label)
    s.graphicalProperties.line.solidFill = color; s.graphicalProperties.line.width = 22000; s.smooth = False; s.marker.symbol = 'none'
    if dash:
        s.graphicalProperties.line.dashStyle = dash
line.y_axis.number_format = '#,##0'; line.y_axis.majorGridlines = None; line.y_axis.delete = False
line.x_axis.number_format = 'mmm yy'; line.x_axis.delete = False; line.x_axis.tickLblSkip = 13
line.set_categories(Reference(eng, min_col=1, min_row=R0, max_row=R1))
line.y_axis.title = 'Sign-ups a week'
line += bar
line.title = None; line.style = 2; line.width = 19.5; line.height = 12.4; line.legend.position = 'b'
play.add_chart(line, 'B10')

# ---- Budget
bud = wb.create_sheet('Budget')
bud.sheet_properties.codeName = 'Sheet2'
btn = app_frame(bud, MAG, '="Chapter 8  ·  The next pound"', '="Where should the next pound go?"',
                '="The third biscuit never tastes like the first: each extra £1 on a channel earns a little less. Tap − and + to move money. Keep the total. Beat the 2025 plan."')
btn('G29:J29', f'=IF({BUDMODE}=1,"Back to my plan","Show AP\'s plan")', 'E9EDF3', NAVY)
btn('L29:O29', 'Reset plan', 'E9EDF3', NAVY)
hdr = ['Channel', '2025 £k', 'Your £k', '', '', 'Sign-ups', 'Return per £1']
for j, (h, colspan) in enumerate(zip(hdr, [(2, 6), (7, 9), (10, 12), (13, 13), (14, 14), (15, 16), (17, 18)])):
    pass
layout = [('B', 'F', 'Channel'), ('G', 'I', '2025 £k'), ('J', 'L', 'Your £k'), ('M', 'M', ''), ('N', 'N', ''), ('O', 'P', 'Sign-ups'), ('Q', 'R', 'Return per £1')]
bud.row_dimensions[10].height = 22
for a, b, h in layout:
    bud.merge_cells(f'{a}10:{b}10'); c = bud[f'{a}10']; c.value = h; c.font = F(10, True, 'FFFFFF'); c.fill = FILL(NAVY); c.alignment = CENTER
for i, (nm, spend, roi, apv) in enumerate(channels, start=1):
    r = 10 + i
    bud.row_dimensions[r].height = 24
    bud.merge_cells(f'B{r}:F{r}'); bud[f'B{r}'] = nm; bud[f'B{r}'].font = F(11); bud[f'B{r}'].alignment = LEFT
    bud.merge_cells(f'G{r}:I{r}'); bud[f'G{r}'] = f'={ST}$C${50 + i}'; bud[f'G{r}'].number_format = '#,##0'; bud[f'G{r}'].font = F(11, color=GREY); bud[f'G{r}'].alignment = CENTER
    bud.merge_cells(f'J{r}:L{r}'); bud[f'J{r}'] = f'=IF({BUDMODE}=1,{ST}$E${50 + i},{BUD(i)})'; bud[f'J{r}'].number_format = '#,##0'; bud[f'J{r}'].font = F(11, True, NAVY); bud[f'J{r}'].alignment = CENTER
    bud[f'M{r}'] = '−'; bud[f'N{r}'] = '+'
    for cc in 'MN':
        bud[f'{cc}{r}'].font = F(13, True, 'FFFFFF'); bud[f'{cc}{r}'].fill = FILL(NAVY); bud[f'{cc}{r}'].alignment = CENTER
    bud.merge_cells(f'O{r}:P{r}'); bud[f'O{r}'] = f'=IF(J{r}<=0,0,({ST}$D${50 + i}*{ST}$C${50 + i}*1000/90)/SQRT({ST}$C${50 + i})*SQRT(J{r}))'
    bud[f'O{r}'].number_format = '#,##0'; bud[f'O{r}'].font = F(11); bud[f'O{r}'].alignment = CENTER
    bud.merge_cells(f'Q{r}:R{r}'); bud[f'Q{r}'] = f'=IF(J{r}<=0,0,O{r}*{POUND}/1000/J{r})'
    bud[f'Q{r}'].number_format = '"£"0.00'; bud[f'Q{r}'].font = F(11); bud[f'Q{r}'].alignment = CENTER
    for cc in 'BGJOQ':
        bud[f'{cc}{r}'].border = Border(bottom=thin)
bud.conditional_formatting.add('J11:J23', DataBarRule(start_type='num', start_value=0, end_type='num', end_value=1600, color=CYAN))
r = 24
bud.row_dimensions[r].height = 22
bud.merge_cells(f'B{r}:F{r}'); bud[f'B{r}'] = 'Total'; bud[f'B{r}'].font = F(11, True)
bud.merge_cells(f'G{r}:I{r}'); bud[f'G{r}'] = '=SUM(G11:G23)'; bud[f'G{r}'].number_format = '#,##0'; bud[f'G{r}'].font = F(11, True); bud[f'G{r}'].alignment = CENTER
bud.merge_cells(f'J{r}:L{r}'); bud[f'J{r}'] = '=SUM(J11:J23)'; bud[f'J{r}'].number_format = '#,##0'; bud[f'J{r}'].font = F(11, True, NAVY); bud[f'J{r}'].alignment = CENTER
bud.merge_cells(f'O{r}:P{r}'); bud[f'O{r}'] = '=SUM(O11:O23)'; bud[f'O{r}'].number_format = '#,##0'; bud[f'O{r}'].font = F(11, True); bud[f'O{r}'].alignment = CENTER
bud.merge_cells(f'Q{r}:R{r}'); bud[f'Q{r}'] = f'=O24*{POUND}/1000/J24'; bud[f'Q{r}'].number_format = '"£"0.00'; bud[f'Q{r}'].font = F(11, True); bud[f'Q{r}'].alignment = CENTER
bud['B26'] = ("Each channel's curve is set so the 2025 plan gives exactly Wave 1's return. Double a channel's spend and you get 41% more sign-ups, not double. "
              "AP's real curves differ by channel, so their plan gained 12,300 sign-ups with them and about 600 here. The direction is the lesson.")
bud['B26'].font = F(8, color=GREY, italic=True); bud['B26'].alignment = TOP; bud.merge_cells('B26:R27')
bud.row_dimensions[26].height = 22; bud.row_dimensions[27].height = 22
# right panel
bud['T10'] = 'Extra sign-ups against the 2025 plan'; bud['T10'].font = F(9, color=GREY)
bud['T11'] = '=O24-SUMPRODUCT(State!$D$51:$D$63,State!$C$51:$C$63)*1000/90'; bud['T11'].number_format = '+#,##0;-#,##0'; bud['T11'].font = F(24, True, NAVY)
bud['T12'] = f'=IF(ABS(J24-G24)<=10,"Same total as 2025. Good.",IF(J24>G24,"Over budget by £"&TEXT(J24-G24,"#,##0")&"k. Take some back.","Under budget by £"&TEXT(G24-J24,"#,##0")&"k. Spend it."))'
bud['T12'].font = F(11); bud['T12'].alignment = LEFT
bud.conditional_formatting.add('T12', FormulaRule(formula=['ISNUMBER(SEARCH("Good",T12))'], font=Font(name='Arial', size=11, color=GREEN)))
bud.conditional_formatting.add('T12', FormulaRule(formula=['NOT(ISNUMBER(SEARCH("Good",T12)))'], font=Font(name='Arial', size=11, color=AMBER)))
bud['T13'] = '=IF(ABS(J24-G24)>10,"Fix the total first.",IF(T11>=1000,"✓ You found "&TEXT(T11,"#,##0")&" extra sign-ups on the same money. The most this table allows is about 2,000.","Move money from the weakest channel to the strongest. Find 1,000 more."))'
bud['T13'].font = F(12); bud['T13'].alignment = TOP; bud.merge_cells('T13:Z15')
bud.conditional_formatting.add('T13', FormulaRule(formula=['LEFT(T13,1)="✓"'], fill=FILL(GREENFILL), font=Font(name='Arial', size=12, color=GREEN)))
bud.conditional_formatting.add('T13', FormulaRule(formula=['LEFT(T13,1)<>"✓"'], fill=FILL(AMBERFILL), font=Font(name='Arial', size=12, color=AMBER)))
bud['T16'] = '=IF(T11>=1000,"Economists call this: marginal return","")'; bud['T16'].font = F(10, True, STEELBLUE)
bud['T18'] = 'The average £1 returned £1.66. The last £1 returned 82p. Should we add £1m to the same plan?'
bud['T18'].font = F(12, True, NAVY); bud['T18'].alignment = TOP; bud.merge_cells('T18:Z19')


def quiz_block(ws, first_row, q_index, opts, tick, nudge, badge):
    rows = [first_row, first_row + 1, first_row + 2]
    for r, (k, txt) in zip(rows, enumerate(opts, start=1)):
        ws[f'T{r}'] = f'{chr(64 + k)}.   {txt}'; ws[f'T{r}'].font = F(12); ws[f'T{r}'].alignment = LEFT
        ws[f'T{r}'].fill = FILL(PALE); ws[f'T{r}'].border = BOX
        ws.row_dimensions[r].height = 24
        ws.conditional_formatting.add(f'T{r}', FormulaRule(formula=[f'{QCHOSEN(q_index)}={k}'], fill=FILL(LIGHTBLUE), font=Font(name='Arial', size=12, bold=True, color=NAVY)))
    fr = first_row + 3
    ws[f'T{fr}'] = f'=IF({QCHOSEN(q_index)}=0,"Pick one.",IF({QCHOSEN(q_index)}={ST}$C${16 + q_index},"✓ {tick}","✗ Not quite. {nudge}"))'
    ws[f'T{fr}'].font = F(12); ws[f'T{fr}'].alignment = TOP; ws.merge_cells(f'T{fr}:Z{fr + 1}')
    ws.row_dimensions[fr].height = 24; ws.row_dimensions[fr + 1].height = 24
    ws.conditional_formatting.add(f'T{fr}', FormulaRule(formula=[f'LEFT(T{fr},1)="✓"'], fill=FILL(GREENFILL), font=Font(name='Arial', size=12, color=GREEN)))
    ws.conditional_formatting.add(f'T{fr}', FormulaRule(formula=[f'LEFT(T{fr},1)="✗"'], fill=FILL(AMBERFILL), font=Font(name='Arial', size=12, color=AMBER)))
    ws[f'T{fr + 2}'] = f'=IF(LEFT(T{fr},1)="✓","Economists call this: {badge}","")'; ws[f'T{fr + 2}'].font = F(10, True, STEELBLUE)
    return rows


quiz_block(bud, 20, 8, ['Not yet: move money first', 'Yes, every £1 returns £1.66', 'No, cut everything'],
           'Moving money between channels was worth about £1.1m before adding a single pound.',
           'The average pound looks great. The next pound does not.', 'Marginal return')

# ---- Long game
lg = wb.create_sheet('Long game')
lg.sheet_properties.codeName = 'Sheet3'
btn = app_frame(lg, CYAN, '="Chapter 9  ·  The long game"', '="Some of today\'s sign-ups come from adverts people saw years ago"',
                '="The pink line is brand consideration: the share of people who would consider supporting us. It moves slowly. Nearly a fifth of it was built by past adverts."')
stats = [('18.8%', 'of today\'s consideration comes from past adverts. That stock is our brand equity.'),
         ('+0.4 points', 'of consideration for every extra £1m of long-term media, against a drift of 0.25 points a year.'),
         ('£2.01', 'is what £1 of Race for Life adverts returns over three years, up from £1.66 in the short term.')]
for k, (big, small) in enumerate(stats):
    r = 10 + k * 3
    lg[f'T{r}'] = big; lg[f'T{r}'].font = F(22, True, MAG); lg.row_dimensions[r].height = 30
    lg[f'T{r + 1}'] = small; lg[f'T{r + 1}'].font = F(11); lg[f'T{r + 1}'].alignment = TOP; lg.merge_cells(f'T{r + 1}:Z{r + 1}'); lg.row_dimensions[r + 1].height = 34
    lg.row_dimensions[r + 2].height = 8
quiz_block(lg, 20, 9, ['£2.01', '£1.66', '82p'],
           'Counting the long-term effect takes Race for Life adverts from £1.66 to £2.01 per pound.',
           'Look at the third figure above.', 'Long-term multiplier')
lg['T19'] = 'Counting the long-term effect, what does £1 of Race for Life adverts return?'; lg['T19'].font = F(12, True, NAVY); lg['T19'].alignment = TOP
lg.merge_cells('T19:Z19'); lg.row_dimensions[19].height = 34
cb = BarChart(); cb.type = 'col'; cb.style = 2; cb.width = 19.5; cb.height = 12.4; cb.legend.position = 'b'
cb.y_axis.majorGridlines = None; cb.y_axis.number_format = '#,##0'; cb.y_axis.title = 'All CRUK media £k a week'
cb.x_axis.number_format = 'mmm yy'; cb.x_axis.delete = False; cb.y_axis.delete = False; cb.x_axis.tickLblSkip = 13; cb.gapWidth = 30
cb.add_data(Reference(eng, min_col=C['All CRUK media £k'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
cb.series[0].graphicalProperties.solidFill = 'D9D9D9'; cb.series[0].tx = SeriesLabel(v='All CRUK media £k')
cb.set_categories(Reference(eng, min_col=1, min_row=R0, max_row=R1))
cl = LineChart()
cl.add_data(Reference(eng, min_col=C['Brand consideration'], min_row=R0 - 1, max_row=R1), titles_from_data=True)
cl.series[0].tx = SeriesLabel(v='Brand consideration'); cl.series[0].graphicalProperties.line.solidFill = MAG
cl.series[0].graphicalProperties.line.width = 22000; cl.series[0].marker.symbol = 'none'; cl.series[0].smooth = False
cl.y_axis.axId = 200; cl.y_axis.crosses = 'max'; cl.y_axis.number_format = '0%'; cl.y_axis.majorGridlines = None; cl.y_axis.delete = False
cl.y_axis.scaling.min = 0.4; cl.y_axis.scaling.max = 0.65; cl.y_axis.title = 'Consideration'
cb += cl
lg.add_chart(cb, 'B10')
lg['B9'] = 'All CRUK media spend each week (grey) and brand consideration (pink)'; lg['B9'].font = F(9, color=GREY)

# ---- Giving
gv = wb.create_sheet('Giving')
gv.sheet_properties.codeName = 'Sheet4'
btn = app_frame(gv, CYAN, '="Chapter 10  ·  Committed Giving"', '="Committed Giving looks like 9p per £1 because the model only counts the first gift"',
                '="The bars show why new monthly givers fell 23% in 2025. Inflation and the advertising cut each took about a quarter. One we chose, one we did not."')
stats = [('73%', 'of new Committed Givers over three years came from marketing, direct and halo together.'),
         ('904 and 307', 'new givers lost to inflation and to the Committed Giving cut. The two biggest causes of the 23% fall.'),
         ('9p per £1', 'counts only the first gift of brand new givers. Not the years that follow, not the churn adverts prevent.')]
for k, (big, small) in enumerate(stats):
    r = 10 + k * 3
    gv[f'T{r}'] = big; gv[f'T{r}'].font = F(22, True, MAG); gv.row_dimensions[r].height = 30
    gv[f'T{r + 1}'] = small; gv[f'T{r + 1}'].font = F(11); gv[f'T{r + 1}'].alignment = TOP; gv.merge_cells(f'T{r + 1}:Z{r + 1}'); gv.row_dimensions[r + 1].height = 34
    gv.row_dimensions[r + 2].height = 8
gv['T19'] = 'Why does Committed Giving show only 9p per £1?'; gv['T19'].font = F(12, True, NAVY); gv['T19'].alignment = TOP
gv.merge_cells('T19:Z19'); gv.row_dimensions[19].height = 34
quiz_block(gv, 20, 10, ['It counts only the first gift of new givers', 'The adverts do not work', 'Inflation'],
           'The short-term return is a floor. The years of giving that follow are not in the number yet.',
           'Think about what a monthly giver is worth over years, not months.', 'Base driver and lifetime value')
waterfall = [('Inflation', -904), ('Committed Giving adverts cut', -307), ('Brand adverts cut', -107), ('Race for Life adverts cut', -102),
             ('Consideration slipped', -29), ('Legacy adverts up', 41), ('Less interest in Macmillan', 46), ('Health campaign halo', 91),
             ('Social Challenges adverts up', 337)]
gv['AC10'] = 'Cause'; gv['AD10'] = 'Change in new givers'
for k, (lab, v) in enumerate(waterfall):
    gv[f'AC{11 + k}'] = lab; gv[f'AD{11 + k}'] = v
gv.column_dimensions['AC'].width = 2; gv.column_dimensions['AD'].width = 2
for k in range(10, 20):
    gv[f'AC{k}'].font = F(1, color='FFFFFF'); gv[f'AD{k}'].font = F(1, color='FFFFFF')
wf = BarChart(); wf.type = 'bar'; wf.style = 2; wf.width = 19.5; wf.height = 12.4; wf.legend = None
wf.y_axis.majorGridlines = None; wf.y_axis.number_format = '#,##0'; wf.x_axis.delete = False; wf.y_axis.delete = False
wf.add_data(Reference(gv, min_col=30, min_row=10, max_row=19), titles_from_data=True)
wf.set_categories(Reference(gv, min_col=29, min_row=11, max_row=19))
wf.series[0].graphicalProperties.solidFill = MAG; wf.series[0].invertIfNegative = False
wf.x_axis.tickLblPos = 'low'
gv.add_chart(wf, 'B10')
gv['B9'] = 'Change in new Committed Givers, FY24/25 to FY25/26 (April to October), by cause. Source: Wave 1'; gv['B9'].font = F(9, color=GREY)

# ---- Finish
fin = wb.create_sheet('Finish')
fin.sheet_properties.codeName = 'Sheet5'
btn = app_frame(fin, NAVY, '="Finish"', f'="You built the model. "&{STARS}&" stars out of 10."',
                '="Ten sentences you can now say in the room, each one true to Wave 1. Tap Start again to play with someone else."', next_label='Start again  ↺')
room = [
    ('Econometrics', 'Econometrics is a detective reading three years of weekly receipts to work out which pounds raised the sign-ups.'),
    ('Baseline', 'Race for Life has a small base of a few hundred sign-ups a week. Nearly every other sign-up was caused by something we did.'),
    ('Pull-forward', 'Sale weeks drove a fifth of sign-ups. The dip after each sale is real, and the model counts it.'),
    ('Adstock', 'Our adverts keep raising sign-ups for weeks after they stop. Judge a campaign in week one and you miss most of what it earned.'),
    ('Halo', 'Committed Giving adverts were quietly filling Race for Life too. When we cut them, Race for Life felt it.'),
    ('Multicollinearity', 'The group number is firm. The channel split is softer. Quote the range when a decision is close.'),
    ('Model tests', 'Wave 1 explains 94% of the weekly pattern, predicts weeks it never saw and left 1.4% of the 2025 fall unexplained.'),
    ('Marginal return', 'The average pound returned £1.66 but the last pound returned 82p. Moving money is worth £1.1m before adding a single pound.'),
    ('Brand equity', 'Nearly a fifth of the people who consider us do so because of adverts they saw years ago. The short-term return is the floor.'),
    ('Lifetime value', 'Committed Giving shows 9p per pound because the model only counts the first gift of brand new givers.'),
]
for k, (word, sent) in enumerate(room):
    r = 9 + k * 2
    fin.row_dimensions[r].height = 30; fin.row_dimensions[r + 1].height = 4
    fin[f'B{r}'] = f'=IF({ST}$B${16 + k + 1}={ST}$C${16 + k + 1},"★","☆")'; fin[f'B{r}'].font = F(16, color=MAG); fin[f'B{r}'].alignment = CENTER
    fin[f'C{r}'] = word; fin[f'C{r}'].font = F(11, True, NAVY); fin.merge_cells(f'C{r}:F{r}'); fin[f'C{r}'].alignment = LEFT
    fin[f'G{r}'] = sent; fin[f'G{r}'].font = F(11); fin.merge_cells(f'G{r}:Z{r}'); fin[f'G{r}'].alignment = LEFT
fin['B29'].value = '◀  Back'
fin.row_dimensions[29].height = 34

# ---- Words (kept short)
wsg = wb.create_sheet('Words')
wsg.sheet_properties.tabColor = GREY
wsg.sheet_view.showGridLines = False
for c, w in {'A': 2, 'B': 24, 'C': 44, 'D': 44, 'E': 44}.items():
    wsg.column_dimensions[c].width = w
wsg['B2'] = 'Words: every term, one rung at a time'; wsg['B2'].font = F(18, True, NAVY)
wsg['B3'] = 'The everyday picture first, then what it meant in Wave 1, then the sentence to say out loud.'; wsg['B3'].font = F(11, italic=True, color=GREY)
for j, h in enumerate(['The word', 'The everyday picture', 'What it meant in Wave 1', 'Say this in the room']):
    c = wsg.cell(row=5, column=2 + j, value=h); c.font = F(11, True, 'FFFFFF'); c.fill = FILL(NAVY); c.alignment = CENTER; c.border = BOX
glossary = [
    ('Econometrics', 'A detective with three years of weekly till receipts, working out which adverts raised the money.', 'AP read 148 weeks of sign-ups against everything that happened each week and measured what each thing was worth.', 'Econometrics is a detective reading three years of weekly receipts to work out which pounds raised the sign-ups.'),
    ('Baseline', 'The customers who walk into a shop because it is there and open.', 'A few hundred a week, from people who already consider us and from our emails. AP say the true base is close to zero.', 'The base is what we would get with the adverts off. Everything above it is what we caused.'),
    ('Driver', 'Anything that moves sign-ups up or down: an advert, a sale, the price, the time of year, inflation.', 'AP tested 18 marketing drivers and six others and kept those with a reliable link. Weather, PR and outages had none.', 'A driver is anything the model found that moves sign-ups.'),
    ('Negative driver', 'A brake. A dearer ticket means fewer people buy it.', 'Price and the shut autumn window pull sign-ups down. A 12% lower fee in 2025 added about 2,500 sign-ups.', 'Price pulls sign-ups down, and the model measures the brake as well as the accelerator.'),
    ('Pull-forward', 'A January sale. People who were going to buy in February buy in January instead, so February is quiet.', 'Fewer sign-ups in the weeks after a discount window closed. AP added this so sales were not over-credited.', 'Discounts bring sign-ups forward as well as adding new ones. The dip afterwards is real.'),
    ('Adstock', 'A kettle. Switch it off and the water stays hot for a while.', 'Each block has a memory. The short-term effect of Race for Life adverts lasts up to six months.', 'Our adverts keep raising sign-ups for weeks after they stop, and the model measures how fast that fades.'),
    ('ROI, return per £1', 'Income earned per pound of adverts. £1.66 back for every £1 in.', 'Sponsorship income only, about £90 per sign-up. Not entry fees, not Gift Aid, only new sign-ups.', 'Every £1 of Race for Life adverts returned £1.66 of sponsorship income in 2025.'),
    ('Halo', 'Rain on one field waters the next field too.', 'Committed Giving adverts were 13% of Race for Life sign-ups in 2023 and 5% in 2025. Measured, not assumed.', 'Committed Giving adverts were quietly filling Race for Life too. That is the halo.'),
    ('Multicollinearity', 'Two people singing the same note. You hear the sound but cannot tell who is louder.', 'Most Race for Life channels air in the same weeks. AP measure the group first, then split it.', 'When channels air together the model is sure about the group and less sure about each channel.'),
    ('Confidence interval', 'The range the true answer probably sits in. A forecast says 15 to 19 degrees, not 17.', 'The Committed Giving halo is 2 to 4 sign-ups per £1,000. Quote the range when a decision is close.', 'Every result is a range, and the range is the honest answer.'),
    ('R squared', 'How much of the weekly up-and-down the model reproduces. 100% would be a perfect copy.', 'Wave 1 explains 94%. AP\'s guide is above 85%.', 'The model explains 94% of the weekly pattern in sign-ups.'),
    ('MAPE', 'On an average week, how far the model is from the real number, as a percentage.', '28% on weeks with more than 2,000 sign-ups. Quiet weeks are left out because tiny weeks make percentages jump.', 'On an average season week the model is within about a quarter of the real number.'),
    ('Holdout test', 'Hide the last few months, fit the model on the rest, then see if it predicts what you hid.', 'The strongest proof a model works. AP use holdouts, bootstrapping and business sense checks together.', 'A good model predicts weeks it was never shown.'),
    ('Unexplained', 'The part of a year-on-year change no driver accounts for. Every model has some.', '1.4% of sign-ups in the 2024 to 2025 comparison, well inside AP\'s tolerance.', 'The model left 1.4% unexplained, which is small.'),
    ('Marginal return', 'The third biscuit never tastes like the first. The return on the very last pound.', 'In 2025 the last pound of Race for Life adverts returned 82p while the average pound returned £1.66.', 'The average pound returned £1.66 but the last pound returned 82p, so adding budget needs care.'),
    ('Brand equity', 'Reputation. It keeps working after you stop paying for it, and drains slowly if you stop.', '18.8% of today\'s consideration traces back to past media. It slips 0.25 points a year. £1m lifts it 0.4 points.', 'Nearly a fifth of the people who consider us do so because of adverts they saw years ago.'),
    ('Long-term multiplier', 'The slow second job of an advert: making more people consider us at all.', '1.21 for Race for Life adverts (£1.66 becomes £2.01) and 1.04 for Committed Giving.', 'Counting the long-term effect takes Race for Life adverts from £1.66 to £2.01 per pound.'),
    ('Base driver', 'The weather for a farmer. It moves the harvest and nobody on the farm controls it.', 'Inflation cost Committed Giving about 900 new givers in FY25/26, as much as the advertising cut.', 'Inflation and the advertising cut each cost about a quarter of new givers. One we chose, one we did not.'),
]
for k, row in enumerate(glossary):
    r = 6 + k
    for colx, txt in zip('BCDE', row):
        c = wsg[f'{colx}{r}']; c.value = txt; c.font = F(11, colx == 'B'); c.alignment = TOP; c.border = BOX
    wsg.row_dimensions[r].height = 15 * max(2, math.ceil(max(len(x) for x in row[1:]) / 44)) + 6

# ---- Room notes (hidden)
rn = wb.create_sheet('Room notes')
rn.column_dimensions['B'].width = 26; rn.column_dimensions['C'].width = 90
rn['B2'] = 'Room notes: running the game with a team in 45 minutes'; rn['B2'].font = F(16, True, NAVY)
rn['B3'] = 'For the person leading. Unhide from the sheet list. The answers are here, so do not share the screen with this tab open.'; rn['B3'].font = F(10, italic=True)
notes = [
    ('0 to 5', 'Open Play on the big screen. Chapter 1, one person tapping. Ask the room what made the spikes before anyone taps.'),
    ('5 to 15', 'Chapters 2 and 3 together. Ask for a show of hands before each number: higher or lower than a thousand? The base surprises people. Let it.'),
    ('15 to 25', 'Chapter 4. Ask the room to guess the memory before you touch it. Most say 20%. Show what 80% does. Name adstock only after the line fits.'),
    ('25 to 32', 'Chapter 5. Add the halo, then read the 42% step out loud and stop. Let the room say what it means for next year.'),
    ('32 to 38', 'Chapter 6 in pairs on laptops. The point is the catch, not the numbers. Ask: which channel are you least sure of, and why?'),
    ('38 to 45', 'Chapter 8 as a race: first to 1,000 extra sign-ups on the same money. Close on average versus last pound. Chapters 9 and 10 are homework.'),
    ('Answer key', f'Base {ans_base:,.0f}. Brake {ans_brake:,.0f}. Sale {ans_sale:,.0f}. Adverts {ans_rfl:.1f} per £1,000, memory {mem_rfl}%. Halo {ans_halo:.2f} per £1,000. Channels: ' + ', '.join(f'{g[0]} {k:.1f}' for g, k in zip(groups, ch_ans)) + '. Every quiz answer is A.'),
    ('If asked "is this just correlation?"', 'Partly, and AP say so. Three things make it more: it controls for everything at once, it is tested on weeks it never saw, and channels that stopped (Race for Life TV after 2024) let it see what happens without them.'),
    ('If asked "the halo felt too big"', 'It did, and AP rebuilt it. They now measure Committed Giving as one block. The halo fell and is still 5% of 2025 sign-ups and 9,139 of the 2025 decline.'),
    ('Budget chapter', 'Best possible on the game\'s curves is about 2,000 extra: cut posters and telemarketing, add DRTV and online video. AP\'s own plan shows the same direction.'),
]
for k, (a, b) in enumerate(notes):
    r = 5 + k
    rn[f'B{r}'] = a; rn[f'B{r}'].font = F(11, True); rn[f'B{r}'].alignment = TOP
    rn[f'C{r}'] = b; rn[f'C{r}'].font = F(11); rn[f'C{r}'].alignment = TOP
    rn.row_dimensions[r].height = 15 * max(2, math.ceil(len(b) / 90)) + 4
rn.sheet_state = 'hidden'

# ---- order and codenames
order = ['Play', 'Budget', 'Long game', 'Giving', 'Finish', 'Words', ENG, STEPS, STATE, 'Room notes', ANS]
wb._sheets = [wb[n] for n in order]
for k, name in enumerate(order, start=1):
    wb[name].sheet_properties.codeName = f'Sheet{k}'
wb.code_name = 'ThisWorkbook'
wb.active = 0
for ws in wb.worksheets:
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.orientation = 'landscape'
    ws.oddHeader.left.text = 'CRUK CONFIDENTIAL, INTERNAL USE ONLY'
    ws.oddFooter.left.text = 'The Econometrics Game v3'
    ws.oddFooter.right.text = 'Andrew Rajanathan, CRUK M&&D'
    if ws.title in ('Play', 'Budget', 'Long game', 'Giving', 'Finish'):
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 1
from openpyxl.workbook.properties import CalcProperties
wb.calculation = CalcProperties(fullCalcOnLoad=True)
wb.save(OUT + '.xlsx')
print('saved', OUT + '.xlsx')
